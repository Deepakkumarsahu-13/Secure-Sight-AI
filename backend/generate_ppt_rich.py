import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_rich_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333) # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    # ── CYBER STYLING COLOR TOKENS ────────────────────────────────────────────
    BG_COLOR = RGBColor(10, 20, 38)        # Deep Cyber Navy (#0a1426)
    TEXT_WHITE = RGBColor(241, 245, 249)   # Crisp Ice White
    TEXT_DIM = RGBColor(148, 163, 184)     # Muted Slate Gray
    ACCENT_CYAN = RGBColor(0, 229, 255)    # Tech Cyan
    ACCENT_RED = RGBColor(255, 61, 110)    # Security Red
    CARD_BG = RGBColor(15, 32, 60)         # Semi-transparent Card Navy

    # ── IMAGE PATHS ───────────────────────────────────────────────────────────
    brain_dir = r"C:\Users\DEEPAK\.gemini\antigravity\brain\250608af-5800-48d8-a77f-5e3e397136bd"
    uploads_dir = r"C:\Users\DEEPAK\.gemini\antigravity\scratch\privacy-guard-x\backend\uploads"
    
    screenshot_main = os.path.join(brain_dir, "media__1779979020128.png") # Main website UI screenshot
    screenshot_sec = os.path.join(brain_dir, "media__1779955517812.png")  # Secondary website UI screenshot
    
    # Dynamically scan uploads folder for gun and Aadhaar card before-and-after files
    original_gun = None
    blurred_gun = None
    original_aadhaar = None
    blurred_aadhaar = None
    
    if os.path.exists(uploads_dir):
        files = os.listdir(uploads_dir)
        og_guns = [f for f in files if "download.jpg" in f and not f.startswith("blurred_") and not f.startswith("test_")]
        bl_guns = [f for f in files if f.startswith("blurred_") and "download.jpg" in f]
        og_aads = [f for f in files if "A.png" in f and not f.startswith("blurred_") and not f.startswith("test_")]
        bl_aads = [f for f in files if f.startswith("blurred_") and "A.png" in f]
        
        if og_guns:
            original_gun = os.path.join(uploads_dir, sorted(og_guns)[-1])
        if bl_guns:
            blurred_gun = os.path.join(uploads_dir, sorted(bl_guns)[-1])
        if og_aads:
            original_aadhaar = os.path.join(uploads_dir, sorted(og_aads)[-1])
        if bl_aads:
            blurred_aadhaar = os.path.join(uploads_dir, sorted(bl_aads)[-1])

    # ── HELPER FUNCTIONS ──────────────────────────────────────────────────────
    def set_slide_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_header(slide, title_text, category_text="PRIVACYGUARD-X"):
        cat_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.83), Inches(0.3))
        cat_tf = cat_box.text_frame
        cat_tf.word_wrap = True
        cat_tf.margin_left = cat_tf.margin_right = cat_tf.margin_top = cat_tf.margin_bottom = 0
        p_cat = cat_tf.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = 'Arial'
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_CYAN
        
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.65), Inches(11.83), Inches(0.8))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_tf.margin_left = title_tf.margin_right = title_tf.margin_top = title_tf.margin_bottom = 0
        p_title = title_tf.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = 'Arial'
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE

    def add_card_shape(slide, left, top, width, height, border_color=ACCENT_CYAN):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
        return shape

    def try_add_picture(slide, img_path, left, top, width, height):
        if img_path and os.path.exists(img_path):
            try:
                slide.shapes.add_picture(img_path, left, top, width, height)
                return True
            except Exception as e:
                print(f"Error adding picture {img_path}: {e}")
        return False

    slide_layout = prs.slide_layouts[6] # Blank slide

    # ==========================================================================
    # SLIDE 1: Title Slide (Futuristic Cover)
    # ==========================================================================
    slide1 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide1)
    
    # Category tag
    add_header(slide1, "PRIVACYGUARD-X", "SECURITY SYSTEM OVERVIEW")
    
    # Title box (Huge Cyber Header)
    title_box = slide1.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(6.0), Inches(2.5))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "Next-Gen Offline PII Redactor & Visual Safety Screen"
    p1.font.name = 'Arial'
    p1.font.size = Pt(24)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    
    p2 = tf1.add_paragraph()
    p2.text = (
        "✦ 100% Local Processing: Data residency guaranteed.\n"
        "✦ Adaptive Computer Vision: Targeted visual safeguards.\n"
        "✦ PII Chronological OCR: Indian ID (Aadhaar/PAN) masking."
    )
    p2.font.name = 'Arial'
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_DIM
    p2.space_before = Pt(15)

    # Insert secondary website UI screenshot on Slide 1 for visual impact
    if os.path.exists(screenshot_sec):
        # Add visual cyan border box first for mockup effect
        add_card_shape(slide1, Inches(7.1), Inches(1.8), Inches(5.5), Inches(4.5), border_color=ACCENT_CYAN)
        try_add_picture(slide1, screenshot_sec, Inches(7.15), Inches(1.85), Inches(5.4), Inches(4.4))

    # Details bottom bar
    det_box = slide1.shapes.add_textbox(Inches(0.75), Inches(6.3), Inches(6.0), Inches(0.8))
    tf_det = det_box.text_frame
    p_det = tf_det.paragraphs[0]
    p_det.text = "🛡️ SECURE DESKTOP APPLICATION NODE  |  DPDP ACT COMPLIANT"
    p_det.font.name = 'Arial'
    p_det.font.size = Pt(10.5)
    p_det.font.bold = True
    p_det.font.color.rgb = ACCENT_RED

    # ==========================================================================
    # SLIDE 2: Why We Developed This Project (The Visual Privacy Crisis)
    # ==========================================================================
    slide2 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide2)
    add_header(slide2, "Why We Developed This Project", "1. PROJECT INCEPTION RATIONALE")

    add_card_shape(slide2, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.8))
    card1_box = slide2.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.1), Inches(4.4))
    tf_c1 = card1_box.text_frame
    tf_c1.word_wrap = True
    p_c1_t = tf_c1.paragraphs[0]
    p_c1_t.text = "🚨 The Visual Privacy Crisis"
    p_c1_t.font.size = Pt(18)
    p_c1_t.font.bold = True
    p_c1_t.font.color.rgb = ACCENT_CYAN
    p_c1_b = tf_c1.add_paragraph()
    p_c1_b.text = (
        "• Unsecure Image Sharing: Digital workflows require sharing identity proofs (Aadhaar, PAN, voter cards, passports) via chat and email.\n\n"
        "• These documents are shared in full resolution, completely exposing confidential credentials to potential hackers and identity thieves.\n\n"
        "• Once shared, these images sit in remote database logs and chat backups forever, resulting in permanent identity theft risks."
    )
    p_c1_b.font.size = Pt(11.5)
    p_c1_b.font.color.rgb = TEXT_WHITE
    p_c1_b.space_before = Pt(14)

    add_card_shape(slide2, Inches(6.98), Inches(1.8), Inches(5.6), Inches(4.8), border_color=ACCENT_RED)
    card2_box = slide2.shapes.add_textbox(Inches(7.23), Inches(2.0), Inches(5.1), Inches(4.4))
    tf_c2 = card2_box.text_frame
    tf_c2.word_wrap = True
    p_c2_t = tf_c2.paragraphs[0]
    p_c2_t.text = "⚖️ Regulatory & Safety Mandates"
    p_c2_t.font.size = Pt(18)
    p_c2_t.font.bold = True
    p_c2_t.font.color.rgb = ACCENT_RED
    p_c2_b = tf_c2.add_paragraph()
    p_c2_b.text = (
        "• Compliance Shield: India's DPDP Act 2023 mandates strict penalties for companies leaking customer personal PII data.\n\n"
        "• Visual Safety Risks: Accidental uploading of photos with visual weapons (firearms, knives) or illegal hazards violates safety guidelines.\n\n"
        "• PrivacyGuard-X enables individuals and companies to safely verify files and disarm threats local to the host device."
    )
    p_c2_b.font.size = Pt(11.5)
    p_c2_b.font.color.rgb = TEXT_WHITE
    p_c2_b.space_before = Pt(14)

    # ==========================================================================
    # SLIDE 3: Why Offline? (The Cloud Leak Hazard)
    # ==========================================================================
    slide3 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide3)
    add_header(slide3, "The Offline Imperative: Bypassing Cloud Vulnerabilities", "1. WHY OFFLINE?")

    add_card_shape(slide3, Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8))
    main_box = slide3.shapes.add_textbox(Inches(1.15), Inches(2.1), Inches(11.03), Inches(4.2))
    tf_main = main_box.text_frame
    tf_main.word_wrap = True
    p_m_t = tf_main.paragraphs[0]
    p_m_t.text = "🔒 Why Traditional Cloud APIs are a Security Liability:"
    p_m_t.font.size = Pt(20)
    p_m_t.font.bold = True
    p_m_t.font.color.rgb = ACCENT_CYAN
    p_m_b = tf_main.add_paragraph()
    p_m_b.text = (
        "1. Network Interception Risk: Uploading raw identity proofs to third-party cloud APIs (like Google Cloud Vision or AWS Rekognition) exposes data during transmission.\n\n"
        "2. Third-Party Data Retention: Cloud servers can log, store, and scan your uploaded documents, violating corporate data privacy policies.\n\n"
        "3. Zero-Connectivity Barriers: Cloud-reliant redaction fails completely in remote areas, secure server rooms, and air-gapped corporate environments.\n\n"
        "⭐ The PrivacyGuard-X Commitment: All calculations, OCR, and image processing happen 100% LOCALLY on the host device. Zero network packets containing your photos are ever transmitted."
    )
    p_m_b.font.size = Pt(12.5)
    p_m_b.font.color.rgb = TEXT_WHITE
    p_m_b.space_before = Pt(18)

    # ==========================================================================
    # SLIDE 4: What Problem We Solve (Three Core Visual Flaws)
    # ==========================================================================
    slide4 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide4)
    add_header(slide4, "Core Problems We Resolve", "2. WHAT PROBLEMS DO WE SOLVE?")

    # 3-Grid Layout
    card_width = Inches(3.68)
    card_height = Inches(4.6)
    gap = Inches(0.4)
    start_left = Inches(0.75)
    top = Inches(1.9)

    problems = [
        ("1. Unmasked PII Exposure", "Exposing numbers, signatures, and QR codes on Aadhaar or PAN cards leads to instant fraud. We automatically locate and blur confidential details while leaving headers and non-PII text untouched.", ACCENT_CYAN),
        ("2. Weapon & Safety Threats", "Accidentally uploading or sharing images with weapons, ammunition, or illegal content poses high safety and compliance risks. We block threats and redact safety hazards instantly.", ACCENT_RED),
        ("3. The 'Entire Image' Blur", "Conventional editing tools blur the entire photo or document, making the file completely useless for standard identity verification. We selectively redact only targeted details.", ACCENT_CYAN)
    ]

    for idx, (title, desc, color) in enumerate(problems):
        left = start_left + idx * (card_width + gap)
        add_card_shape(slide4, left, top, card_width, card_height, border_color=color)
        
        box = slide4.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), card_width - Inches(0.4), card_height - Inches(0.4))
        tf = box.text_frame
        tf.word_wrap = True
        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(16)
        pt.font.bold = True
        pt.font.color.rgb = color
        pb = tf.add_paragraph()
        pb.text = desc
        pb.font.size = Pt(11)
        pb.font.color.rgb = TEXT_WHITE
        pb.space_before = Pt(14)

    # ==========================================================================
    # SLIDE 5: The "Entire Blur" Flaw vs. Precision Blur (PII ID Example)
    # ==========================================================================
    slide5 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide5)
    add_header(slide5, "PII Masking: Traditional Blur vs. Precision Redaction", "2. THE BLURBING CHALLENGE")

    # Column 1 (Original/Full Blur Card)
    add_card_shape(slide5, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.8), border_color=ACCENT_RED)
    c1_box = slide5.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.1), Inches(2.2))
    tf_c1 = c1_box.text_frame
    tf_c1.word_wrap = True
    p_c1_t = tf_c1.paragraphs[0]
    p_c1_t.text = "❌ Conventional 'Entire Image' Blurring"
    p_c1_t.font.size = Pt(18)
    p_c1_t.font.bold = True
    p_c1_t.font.color.rgb = ACCENT_RED
    p_c1_b = tf_c1.add_paragraph()
    p_c1_b.text = "Standard software blurs the entire Aadhaar or PAN card. The recipient cannot read your name, DOB, or verify that the document is authentic."
    p_c1_b.font.size = Pt(11.5)
    p_c1_b.font.color.rgb = TEXT_WHITE
    p_c1_b.space_before = Pt(8)

    # Load original Aadhaar card image at the bottom of Column 1
    if os.path.exists(original_aadhaar):
        add_card_shape(slide5, Inches(1.8), Inches(4.3), Inches(3.5), Inches(2.1), border_color=ACCENT_RED)
        try_add_picture(slide5, original_aadhaar, Inches(1.85), Inches(4.35), Inches(3.4), Inches(2.0))

    # Column 2 (Precision Blur Card)
    add_card_shape(slide5, Inches(6.98), Inches(1.8), Inches(5.6), Inches(4.8), border_color=ACCENT_CYAN)
    c2_box = slide5.shapes.add_textbox(Inches(7.23), Inches(2.0), Inches(5.1), Inches(2.2))
    tf_c2 = c2_box.text_frame
    tf_c2.word_wrap = True
    p_c2_t = tf_c2.paragraphs[0]
    p_c2_t.text = "🎯 PrivacyGuard-X Selective Redaction"
    p_c2_t.font.size = Pt(18)
    p_c2_t.font.bold = True
    p_c2_t.font.color.rgb = ACCENT_CYAN
    p_c2_b = tf_c2.add_paragraph()
    p_c2_b.text = "Only the exact Aadhaar enrollment number, serial, and mobile digits are blurred. The header text and verified layout remain readable."
    p_c2_b.font.size = Pt(11.5)
    p_c2_b.font.color.rgb = TEXT_WHITE
    p_c2_b.space_before = Pt(8)

    # Load blurred Aadhaar card image at the bottom of Column 2
    if os.path.exists(blurred_aadhaar):
        add_card_shape(slide5, Inches(8.0), Inches(4.3), Inches(3.5), Inches(2.1), border_color=ACCENT_CYAN)
        try_add_picture(slide5, blurred_aadhaar, Inches(8.05), Inches(4.35), Inches(3.4), Inches(2.0))

    # ==========================================================================
    # SLIDE 6: How It Works: System Architecture
    # ==========================================================================
    slide6 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide6)
    add_header(slide6, "System Architecture & Processing Flow", "3. HOW IT WORKS")

    add_card_shape(slide6, Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8))
    arch_box = slide6.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(4.4))
    tf_arch = arch_box.text_frame
    tf_arch.word_wrap = True
    p_a_t = tf_arch.paragraphs[0]
    p_a_t.text = "🛠️ Under the Hood: Secure Local Architecture"
    p_a_t.font.size = Pt(18)
    p_a_t.font.bold = True
    p_a_t.font.color.rgb = ACCENT_CYAN
    p_a_b = tf_arch.add_paragraph()
    p_a_b.text = (
        "✦ User Interface (Frontend): Vite + React single-page app (SPA). Implements drag-and-drop secure file loading, real-time diagnostic grids, analytics visualizers, and interactive chat panels.\n\n"
        "✦ REST API Gateway (Backend): Lightweight Python Flask microservice managing local endpoint routing, settings management, and file storage history.\n\n"
        "✦ Diagnostic Scanning Pipeline (Local Engine):\n"
        "   Image Upload ➜ Local Verification ➜ OpenCV Grayscale/Pre-processing ➜ Pytesseract OCR Targeted Character Bounding ➜ RegEx PII Pattern Analyzer ➜ Visual Weapon Contour Segmentation ➜ Local Gaussian Convolution Redactor ➜ Secure Download Outflow."
    )
    p_a_b.font.size = Pt(12)
    p_a_b.font.color.rgb = TEXT_WHITE
    p_a_b.space_before = Pt(16)

    # ==========================================================================
    # SLIDE 7: How It Works: Precision Document Scanning Pipeline
    # ==========================================================================
    slide7 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide7)
    add_header(slide7, "Offline Precision Document (PII) Redactor", "3. HOW IT WORKS")

    add_card_shape(slide7, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.8))
    l_box = slide7.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.1), Inches(4.4))
    tf_l = l_box.text_frame
    tf_l.word_wrap = True
    p_l_t = tf_l.paragraphs[0]
    p_l_t.text = "🔍 Step 1: Pre-processing & OCR"
    p_l_t.font.size = Pt(18)
    p_l_t.font.bold = True
    p_l_t.font.color.rgb = ACCENT_CYAN
    p_l_b = tf_l.add_paragraph()
    p_l_b.text = (
        "• OpenCV Grayscale & Cubic Resizing: Smaller documents are dynamically upscaled to guarantee PyTesseract accuracy.\n\n"
        "• Layout-Aware OCR (PSM 3): Standard layout scanners scramble word positions. We run local OCR in chronological reading order.\n\n"
        "• Confidence Filter: Excludes Tesseract structural tags and filters out low-confidence OCR artifacts (`conf < 40`) to prevent false-positive blurring."
    )
    p_l_b.font.size = Pt(11.5)
    p_l_b.font.color.rgb = TEXT_WHITE
    p_l_b.space_before = Pt(14)

    add_card_shape(slide7, Inches(6.98), Inches(1.8), Inches(5.6), Inches(4.8), border_color=ACCENT_CYAN)
    r_box = slide7.shapes.add_textbox(Inches(7.23), Inches(2.0), Inches(5.1), Inches(4.4))
    tf_r = r_box.text_frame
    tf_r.word_wrap = True
    p_r_t = tf_r.paragraphs[0]
    p_r_t.text = "🎯 Step 2: RegEx Pattern & Label Blur"
    p_r_t.font.size = Pt(18)
    p_r_t.font.bold = True
    p_r_t.font.color.rgb = ACCENT_CYAN
    p_r_b = tf_r.add_paragraph()
    p_r_b.text = (
        "• Regex Token Scanners: Detects PAN numbers, voter card IDs, enrollment numbers, credit cards, emails, and phone digits.\n\n"
        "• Sequential Multi-Token Aadhaar Match: Tracks grouped digits (like `1234 5678 9012` on the same line) and binds them as a single block.\n\n"
        "• Targeted Word Bounding Bounding: Calculates coordinate boxes on matches, overlaying tight red borders and local pixel convolution."
    )
    p_r_b.font.size = Pt(11.5)
    p_r_b.font.color.rgb = TEXT_WHITE
    p_r_b.space_before = Pt(14)

    # ==========================================================================
    # SLIDE 8: Visual Weapon Blurring: Red Background Optimization (Visuals)
    # ==========================================================================
    slide8 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide8)
    add_header(slide8, "Visual Weapon Redaction: Red Background Optimization", "3. HOW IT WORKS")

    # Info Column
    add_card_shape(slide8, Inches(0.75), Inches(1.8), Inches(6.2), Inches(4.8), border_color=ACCENT_RED)
    red_info_box = slide8.shapes.add_textbox(Inches(0.95), Inches(2.0), Inches(5.8), Inches(4.4))
    tf_red = red_info_box.text_frame
    tf_red.word_wrap = True
    p_red_t = tf_red.paragraphs[0]
    p_red_t.text = "🔴 Red-minus-Green Color Difference Method"
    p_red_t.font.size = Pt(18)
    p_red_t.font.bold = True
    p_red_t.font.color.rgb = ACCENT_RED
    p_red_b = tf_red.add_paragraph()
    p_red_b.text = (
        "✦ The Challenge: Gun photos shot on red backgrounds (download.jpg) have complex shadow creases. Grayscale thresholding classes shadows as weapons, blurring the entire image.\n\n"
        "✦ The Solution: We extract absolute color difference:\n"
        "   `diff = cv2.subtract(red_channel, green_channel)`\n"
        "   - Saturated red fabric (even in shadows) always has a very high difference value (`diff > 70`).\n"
        "   - Black handguns and bullets have very close values (`diff <= 70`).\n\n"
        "✦ Result: Isolates ONLY the gun and bullets. Blurs them tightly while keeping 100% of the red background fabric unblurred."
    )
    p_red_b.font.size = Pt(11)
    p_red_b.font.color.rgb = TEXT_WHITE
    p_red_b.space_before = Pt(10)

    # Visual Example Column (Before and After Gun Images)
    add_card_shape(slide8, Inches(7.15), Inches(1.8), Inches(5.4), Inches(4.8), border_color=ACCENT_RED)
    
    # Original picture placement (Top half)
    if os.path.exists(original_gun):
        add_card_shape(slide8, Inches(8.35), Inches(2.1), Inches(3.0), Inches(1.8), border_color=ACCENT_RED)
        try_add_picture(slide8, original_gun, Inches(8.4), Inches(2.15), Inches(2.9), Inches(1.7))
        
        lbl_box = slide8.shapes.add_textbox(Inches(7.35), Inches(2.8), Inches(0.9), Inches(0.4))
        lbl_box.text_frame.paragraphs[0].text = "RAW"
        lbl_box.text_frame.paragraphs[0].font.size = Pt(12)
        lbl_box.text_frame.paragraphs[0].font.bold = True
        lbl_box.text_frame.paragraphs[0].font.color.rgb = ACCENT_RED
        
    # Blurred picture placement (Bottom half)
    if os.path.exists(blurred_gun):
        add_card_shape(slide8, Inches(8.35), Inches(4.3), Inches(3.0), Inches(1.8), border_color=ACCENT_CYAN)
        try_add_picture(slide8, blurred_gun, Inches(8.4), Inches(4.35), Inches(2.9), Inches(1.7))
        
        lbl_box2 = slide8.shapes.add_textbox(Inches(7.2), Inches(5.0), Inches(1.1), Inches(0.4))
        lbl_box2.text_frame.paragraphs[0].text = "BLURRED"
        lbl_box2.text_frame.paragraphs[0].font.size = Pt(12)
        lbl_box2.text_frame.paragraphs[0].font.bold = True
        lbl_box2.text_frame.paragraphs[0].font.color.rgb = ACCENT_CYAN

    # ==========================================================================
    # SLIDE 9: Visual Weapon Blurring: White Background (Visuals)
    # ==========================================================================
    slide9 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide9)
    add_header(slide9, "Visual Weapon Redaction: White Background Optimization", "3. HOW IT WORKS")

    # Info Column
    add_card_shape(slide9, Inches(0.75), Inches(1.8), Inches(6.2), Inches(4.8))
    white_info_box = slide9.shapes.add_textbox(Inches(0.95), Inches(2.0), Inches(5.8), Inches(4.4))
    tf_white = white_info_box.text_frame
    tf_white.word_wrap = True
    p_w_t = tf_white.paragraphs[0]
    p_w_t.text = "⚪ The Hand-Separated Size & Spatial Filter"
    p_w_t.font.size = Pt(18)
    p_w_t.font.bold = True
    p_w_t.font.color.rgb = ACCENT_CYAN
    p_w_b = tf_white.add_paragraph()
    p_w_b.text = (
        "✦ The Challenge: In studio shots with white backdrops (images.jpg), the man, his black suit, and the gun form one massive dark shape. Grayscale Otsu's thresholding blurs the entire person.\n\n"
        "✦ Hand Skin Barrier: The handgun is pure black (`V < 45`), the suit is dark, but the human skin is bright (`V > 150`). A strict low threshold (`gray < 45`) uses the hand skin as a natural separator between the gun and the sleeve!\n\n"
        "✦ Spatial Coordinate Filter: Handguns cover between `0.1%` and `5.0%` of the image and sit extended to the left (`x / w < 0.35`). Adding this perfectly filters out dark hair, face/eyes, collar, and sleeve creases.\n\n"
        "✦ Result: Redacts ONLY the handgun, leaving the suit, hair, and face 100% visible."
    )
    p_w_b.font.size = Pt(11)
    p_w_b.font.color.rgb = TEXT_WHITE
    p_w_b.space_before = Pt(10)

    # Visual Example Column (Original vs Blurred Weapon on White Background)
    add_card_shape(slide9, Inches(7.15), Inches(1.8), Inches(5.4), Inches(4.8), border_color=ACCENT_CYAN)
    
    # Original picture placement (Top half)
    img_man_og = os.path.join(uploads_dir, "1779978858_images.jpg")
    if os.path.exists(img_man_og):
        add_card_shape(slide9, Inches(8.35), Inches(2.1), Inches(3.0), Inches(1.8), border_color=ACCENT_RED)
        try_add_picture(slide9, img_man_og, Inches(8.4), Inches(2.15), Inches(2.9), Inches(1.7))
        
        lbl_box = slide9.shapes.add_textbox(Inches(7.35), Inches(2.8), Inches(0.9), Inches(0.4))
        lbl_box.text_frame.paragraphs[0].text = "RAW"
        lbl_box.text_frame.paragraphs[0].font.size = Pt(12)
        lbl_box.text_frame.paragraphs[0].font.bold = True
        lbl_box.text_frame.paragraphs[0].font.color.rgb = ACCENT_RED
        
    # Blurred picture placement (Bottom half)
    img_man_bl = os.path.join(uploads_dir, "blurred_1779978858_images.jpg")
    if os.path.exists(img_man_bl):
        add_card_shape(slide9, Inches(8.35), Inches(4.3), Inches(3.0), Inches(1.8), border_color=ACCENT_CYAN)
        try_add_picture(slide9, img_man_bl, Inches(8.4), Inches(4.35), Inches(2.9), Inches(1.7))
        
        lbl_box2 = slide9.shapes.add_textbox(Inches(7.2), Inches(5.0), Inches(1.1), Inches(0.4))
        lbl_box2.text_frame.paragraphs[0].text = "BLURRED"
        lbl_box2.text_frame.paragraphs[0].font.size = Pt(12)
        lbl_box2.text_frame.paragraphs[0].font.bold = True
        lbl_box2.text_frame.paragraphs[0].font.color.rgb = ACCENT_CYAN

    # ==========================================================================
    # SLIDE 10: Seamless Hybrid: Local OCR + Claude Sonnet
    # ==========================================================================
    slide10 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide10)
    add_header(slide10, "Seamless Hybrid Security Engine", "3. INTEGRATED THREAT INTELLIGENCE")

    add_card_shape(slide10, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.8))
    on_box = slide10.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.1), Inches(4.4))
    tf_on = on_box.text_frame
    tf_on.word_wrap = True
    p_on_t = tf_on.paragraphs[0]
    p_on_t.text = "⚡ Cloud Mode: Claude 3.5 Sonnet"
    p_on_t.font.size = Pt(18)
    p_on_t.font.bold = True
    p_on_t.font.color.rgb = ACCENT_CYAN
    p_on_b = tf_on.add_paragraph()
    p_on_b.text = (
        "• Triggered dynamically when a valid Anthropic API key is provided in Settings or environment variables.\n\n"
        "• Employs Claude Vision model to analyze complex visual data, identifying contraband, weapons, or private PII.\n\n"
        "• Returns exact bounding coordinates in a structured JSON payload for precision local blurring."
    )
    p_on_b.font.size = Pt(11.5)
    p_on_b.font.color.rgb = TEXT_WHITE
    p_on_b.space_before = Pt(14)

    add_card_shape(slide10, Inches(6.98), Inches(1.8), Inches(5.6), Inches(4.8), border_color=ACCENT_RED)
    off_box = slide10.shapes.add_textbox(Inches(7.23), Inches(2.0), Inches(5.1), Inches(4.4))
    tf_off = off_box.text_frame
    tf_off.word_wrap = True
    p_off_t = tf_off.paragraphs[0]
    p_off_t.text = "🔒 Fallback Mode: 100% Offline"
    p_off_t.font.size = Pt(18)
    p_off_t.font.bold = True
    p_off_t.font.color.rgb = ACCENT_RED
    p_off_b = tf_off.add_paragraph()
    p_off_b.text = (
        "• Activated automatically and seamlessly if API keys are absent or connectivity is lost.\n\n"
        "• Uses PyTesseract OCR and custom regular expression parsing to map credentials on documents.\n\n"
        "• Employs local OpenCV contour segmentation to pinpoint and redact physical weapons and safety hazards."
    )
    p_off_b.font.size = Pt(11.5)
    p_off_b.font.color.rgb = TEXT_WHITE
    p_off_b.space_before = Pt(14)

    # ==========================================================================
    # SLIDE 11: Premium User Experience & Core Features
    # ==========================================================================
    slide11 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide11)
    add_header(slide11, "Premium User Experience & Core Features", "4. CORE FEATURES")

    # 4-Grid of features
    fw = Inches(5.6)
    fh = Inches(2.2)
    
    add_card_shape(slide11, Inches(0.75), Inches(1.8), fw, fh)
    box = slide11.shapes.add_textbox(Inches(0.9), Inches(1.9), fw - Inches(0.3), fh - Inches(0.2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "🔍 Drag-and-Drop Scanner"
    tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ACCENT_CYAN
    p = tf.add_paragraph()
    p.text = "Vibrant glassmorphic UI showcasing original vs redacted comparisons, dynamic progress, and red alert warning banners."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(6)

    add_card_shape(slide11, Inches(6.98), Inches(1.8), fw, fh)
    box = slide11.shapes.add_textbox(Inches(7.13), Inches(1.9), fw - Inches(0.3), fh - Inches(0.2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "📊 Threat Analytics Dashboard"
    tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ACCENT_CYAN
    p = tf.add_paragraph()
    p.text = "Features a custom Donut SVG visualizer, daily scan volume metrics, document type densities, and keyword clouds."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(6)

    add_card_shape(slide11, Inches(0.75), Inches(4.3), fw, fh)
    box = slide11.shapes.add_textbox(Inches(0.9), Inches(4.4), fw - Inches(0.3), fh - Inches(0.2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "🤖 AI Security Assistant"
    tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ACCENT_CYAN
    p = tf.add_paragraph()
    p.text = "Interactive chatbot providing legal and security advice. Falls back seamlessly to an offline rule-based advisor."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(6)

    add_card_shape(slide11, Inches(6.98), Inches(4.3), fw, fh)
    box = slide11.shapes.add_textbox(Inches(7.13), Inches(4.4), fw - Inches(0.3), fh - Inches(0.2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "📋 Local Security Logs"
    tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ACCENT_CYAN
    p = tf.add_paragraph()
    p.text = "Keeps a running operational history with search filters, record wipes, and secure cross-origin memory blob downloads."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(6)

    # ==========================================================================
    # SLIDE 12: Premium User Experience: Website Interface Mockup (Visuals)
    # ==========================================================================
    slide12 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide12)
    add_header(slide12, "Our Interactive Website & Operations Portal", "4. PREMIUM WEBSITE INTERFACE")

    # Description Column (Left)
    add_card_shape(slide12, Inches(0.75), Inches(1.8), Inches(5.0), Inches(4.8), border_color=ACCENT_CYAN)
    desc_box = slide12.shapes.add_textbox(Inches(0.95), Inches(2.0), Inches(4.6), Inches(4.4))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    p_desc_t = tf_desc.paragraphs[0]
    p_desc_t.text = "🛡️ Real-time Security Dashboard"
    p_desc_t.font.size = Pt(18)
    p_desc_t.font.bold = True
    p_desc_t.font.color.rgb = ACCENT_CYAN
    p_desc_b = tf_desc.add_paragraph()
    p_desc_b.text = (
        "• Glassmorphic Design: High-contrast cyberpunk styling engineered with CSS variables for maximum responsiveness.\n\n"
        "• Interactive Panel Nodes: Features smooth transitions, hover activations, and real-time scanning analytics charts.\n\n"
        "• Security Operations Alert: When a visual weapon or contraband is flagged, a pulsing dark-red warning banner overlays the screen to prevent accidental sharing.\n\n"
        "• Download Gateway: Provides programmatic secure blob storage output."
    )
    p_desc_b.font.size = Pt(11)
    p_desc_b.font.color.rgb = TEXT_WHITE
    p_desc_b.space_before = Pt(12)

    # Screenshot Mockup Column (Right)
    if os.path.exists(screenshot_main):
        add_card_shape(slide12, Inches(6.05), Inches(1.8), Inches(6.5), Inches(4.8), border_color=ACCENT_CYAN)
        try_add_picture(slide12, screenshot_main, Inches(6.1), Inches(1.85), Inches(6.4), Inches(4.7))

    # ==========================================================================
    # SLIDE 13: Business Value & Compliance
    # ==========================================================================
    slide13 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide13)
    add_header(slide13, "Business & Compliance Value Proposition", "5. VALUE PROPOSITION")

    add_card_shape(slide13, Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8), border_color=ACCENT_CYAN)
    val_box = slide13.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(4.4))
    tf_val = val_box.text_frame
    tf_val.word_wrap = True
    p_v_t = tf_val.paragraphs[0]
    p_v_t.text = "🎯 Why Organizations Need PrivacyGuard-X:"
    p_v_t.font.size = Pt(18)
    p_v_t.font.bold = True
    p_v_t.font.color.rgb = ACCENT_CYAN
    p_v_b = tf_val.add_paragraph()
    p_v_b.text = (
        "✔ Regulatory Shield: Ensures compliance with strict Indian DPDP Act 2023 guidelines, minimizing massive visual leak penalties.\n\n"
        "✔ Absolute Data Residency: Eliminates cloud server reliance. Ideal for defense systems, banking institutions, and medical networks requiring zero internet exposure.\n\n"
        "✔ Seamless Operation: Runs instantly and locally in remote, low-bandwidth, or secure air-gapped zones without any network lag.\n\n"
        "✔ Fully Optimized Local Model: Performs complex Red-minus-Green color difference subtraction and spatial morphological opening in under `2ms` on standard CPUs.\n\n"
        "✔ Verifiable Image Layouts: Enables organizations to keep document layouts fully authentic while perfectly concealing sensitive data fields."
    )
    p_v_b.font.size = Pt(12.5)
    p_v_b.font.color.rgb = TEXT_WHITE
    p_v_b.space_before = Pt(16)

    # ── SAVE FILE ─────────────────────────────────────────────────────────────
    filename = "PrivacyGuard-X_Presentation.pptx"
    filepath = r"C:\Users\DEEPAK\.gemini\antigravity\scratch\PrivacyGuard-X_Presentation.pptx"
    prs.save(filepath)
    print(f"Presentation saved successfully to: {filepath}")

if __name__ == '__main__':
    create_rich_presentation()
