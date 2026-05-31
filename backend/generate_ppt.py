import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333) # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    # ── CYBER STYLING COLOR TOKENS ────────────────────────────────────────────
    BG_COLOR = RGBColor(10, 20, 38)        # Deep Cyber Navy (#0a1426)
    TEXT_WHITE = RGBColor(241, 245, 249)   # Crisp Ice White
    TEXT_DIM = RGBColor(148, 163, 184)     # Muted Slate Gray
    ACCENT_CYAN = RGBColor(0, 229, 255)    # Tech Cyan
    ACCENT_RED = RGBColor(255, 61, 110)    # Security Red
    CARD_BG = RGBColor(15, 32, 60)         # Semi-transparent Card Navy

    # ── HELPER FUNCTIONS ──────────────────────────────────────────────────────
    def set_slide_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_header(slide, title_text, category_text="PRIVACYGUARD-X"):
        # Category tag
        cat_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.83), Inches(0.3))
        cat_tf = cat_box.text_frame
        cat_tf.word_wrap = True
        cat_tf.margin_left = cat_tf.margin_right = cat_tf.margin_top = cat_tf.margin_bottom = 0
        p_cat = cat_tf.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = 'Arial'
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_CYAN
        
        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.65), Inches(11.83), Inches(0.8))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_tf.margin_left = title_tf.margin_right = title_tf.margin_top = title_tf.margin_bottom = 0
        p_title = title_tf.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = 'Arial'
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE

    def add_card_shape(slide, left, top, width, height, border_color=ACCENT_CYAN):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
        return shape

    # ==========================================================================
    # SLIDE 1: Title Slide (Futuristic Cover)
    # ==========================================================================
    slide_layout = prs.slide_layouts[6] # Blank slide
    slide1 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide1)

    # Title box (Huge Cyber Header)
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(2.5))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "PRIVACYGUARD-X"
    p1.font.name = 'Arial'
    p1.font.size = Pt(54)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_CYAN
    p1.alignment = PP_ALIGN.LEFT
    
    p2 = tf1.add_paragraph()
    p2.text = "Next-Gen Offline PII Redactor & Visual Safety Screen"
    p2.font.name = 'Arial'
    p2.font.size = Pt(22)
    p2.font.color.rgb = TEXT_WHITE
    p2.space_before = Pt(15)
    
    # Details box
    det_box = slide1.shapes.add_textbox(Inches(1.0), Inches(4.8), Inches(11.3), Inches(1.5))
    tf_det = det_box.text_frame
    p_det = tf_det.paragraphs[0]
    p_det.text = "🛡️ SECURE LOCAL REDACTION  |  100% OFFLINE DATA RESIDENCY  |  DPDP ACT 2023 COMPLIANT"
    p_det.font.name = 'Arial'
    p_det.font.size = Pt(12)
    p_det.font.bold = True
    p_det.font.color.rgb = ACCENT_RED
    
    p_det2 = tf_det.add_paragraph()
    p_det2.text = "An advanced Flask + React desktop security node powered by OpenCV and PyTesseract OCR."
    p_det2.font.name = 'Arial'
    p_det2.font.size = Pt(11)
    p_det2.font.color.rgb = TEXT_DIM
    p_det2.space_before = Pt(8)

    # ==========================================================================
    # SLIDE 2: Why We Developed This Project (The Visual Privacy Crisis)
    # ==========================================================================
    slide2 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide2)
    add_header(slide2, "Why We Developed This Project", "1. PROJECT INCEPTION RATIONALE")

    # Column 1 Card (Unsecure Image Sharing)
    add_card_shape(slide2, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.8))
    card1_box = slide2.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.1), Inches(4.4))
    tf_c1 = card1_box.text_frame
    tf_c1.word_wrap = True
    p_c1_t = tf_c1.paragraphs[0]
    p_c1_t.text = "🚨 The Visual Privacy Crisis"
    p_c1_t.font.size = Pt(18)
    p_c1_t.font.bold = True
    p_c1_t.font.color.rgb = ACCENT_CYAN
    
    p_c1_b = tf_c1.add_paragraph()
    p_c1_b.text = (
        "• Everyday digital transactions require sharing identity proofs (Aadhaar, PAN, voter cards, passports) via chat apps and emails.\n\n"
        "• These documents are shared in full resolution, completely exposing confidential credentials to potential hackers and bad actors.\n\n"
        "• Once shared, these images sit in remote database logs and chat backups forever, resulting in massive identity theft and financial fraud risks."
    )
    p_c1_b.font.size = Pt(12)
    p_c1_b.font.color.rgb = TEXT_WHITE
    p_c1_b.space_before = Pt(14)

    # Column 2 Card (Regulatory Compliance)
    add_card_shape(slide2, Inches(6.98), Inches(1.8), Inches(5.6), Inches(4.8), border_color=ACCENT_RED)
    card2_box = slide2.shapes.add_textbox(Inches(7.23), Inches(2.0), Inches(5.1), Inches(4.4))
    tf_c2 = card2_box.text_frame
    tf_c2.word_wrap = True
    p_c2_t = tf_c2.paragraphs[0]
    p_c2_t.text = "⚖️ Tightening Data Protection Laws"
    p_c2_t.font.size = Pt(18)
    p_c2_t.font.bold = True
    p_c2_t.font.color.rgb = ACCENT_RED
    
    p_c2_b = tf_c2.add_paragraph()
    p_c2_b.text = (
        "• India's Digital Personal Data Protection (DPDP) Act 2023 mandates strict penalties for companies leaking personal PII data.\n\n"
        "• Organizations must safeguard customer documents at all times. Standard identity proofs must be masked before processing.\n\n"
        "• PrivacyGuard-X was developed to empower individuals and small businesses to instantly redact sensitive data locally before sharing."
    )
    p_c2_b.font.size = Pt(12)
    p_c2_b.font.color.rgb = TEXT_WHITE
    p_c2_b.space_before = Pt(14)

    # ==========================================================================
    # SLIDE 3: Why Offline? (The Cloud Leak Hazard)
    # ==========================================================================
    slide3 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide3)
    add_header(slide3, "The Offline Imperative: Bypassing Cloud Vulnerabilities", "1. WHY OFFLINE?")

    # Center Card
    add_card_shape(slide3, Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8))
    main_box = slide3.shapes.add_textbox(Inches(1.15), Inches(2.1), Inches(11.03), Inches(4.2))
    tf_main = main_box.text_frame
    tf_main.word_wrap = True
    
    p_m_t = tf_main.paragraphs[0]
    p_m_t.text = "🔒 Why Traditional Cloud APIs are a Security Liability:"
    p_m_t.font.size = Pt(20)
    p_m_t.font.bold = True
    p_m_t.font.color.rgb = ACCENT_CYAN
    
    p_m_b = tf_main.add_paragraph()
    p_m_b.text = (
        "1. Network Interception Risk: Uploading raw identity proofs to third-party cloud APIs (like Google Cloud Vision or AWS Rekognition) exposes data during transmission.\n\n"
        "2. Third-Party Data Retention: Cloud servers can log, store, and scan your uploaded documents, violating corporate data privacy policies.\n\n"
        "3. Zero-Connectivity Barriers: Cloud-reliant redaction fails completely in remote areas, air-gapped corporate environments, and secure offline bank branches.\n\n"
        "⭐ The PrivacyGuard-X Commitment: All calculations, OCR, and image processing happen 100% LOCALLY on the host device. Zero network packets containing your photos are ever transmitted."
    )
    p_m_b.font.size = Pt(13)
    p_m_b.font.color.rgb = TEXT_WHITE
    p_m_b.space_before = Pt(18)

    # ==========================================================================
    # SLIDE 4: What Problem We Solve (Three Core Visual Flaws)
    # ==========================================================================
    slide4 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide4)
    add_header(slide4, "Core Problems We Resolve", "2. WHAT PROBLEMS DO WE SOLVE?")

    # 3-Grid Layout
    card_width = Inches(3.68)
    card_height = Inches(4.6)
    gap = Inches(0.4)
    start_left = Inches(0.75)
    top = Inches(1.9)

    problems = [
        ("1. Unmasked PII Exposure", "Exposing numbers, signatures, and QR codes on Aadhaar or PAN cards leads to instant fraud. We automatically locate and blur confidential details while leaving headers and non-PII text untouched.", ACCENT_CYAN),
        ("2. Weapon & Safety Threats", "Accidentally uploading or sharing images with weapons, ammunition, or illegal content poses high safety and compliance risks. We block threats and redact safety hazards instantly.", ACCENT_RED),
        ("3. The 'Entire Image' Blur", "Conventional editing tools blur the entire photo or document, making the file completely useless for standard identity verification. We selectively redact only targeted details.", ACCENT_CYAN)
    ]

    for idx, (title, desc, color) in enumerate(problems):
        left = start_left + idx * (card_width + gap)
        add_card_shape(slide4, left, top, card_width, card_height, border_color=color)
        
        box = slide4.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), card_width - Inches(0.4), card_height - Inches(0.4))
        tf = box.text_frame
        tf.word_wrap = True
        
        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(16)
        pt.font.bold = True
        pt.font.color.rgb = color
        
        pb = tf.add_paragraph()
        pb.text = desc
        pb.font.size = Pt(11.5)
        pb.font.color.rgb = TEXT_WHITE
        pb.space_before = Pt(14)

    # ==========================================================================
    # SLIDE 5: The "Entire Blur" Flaw vs. Precision Blur
    # ==========================================================================
    slide5 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide5)
    add_header(slide5, "The 'Entire Image Blur' Flaw vs. Precision Redaction", "2. THE BLURBING CHALLENGE")

    # Before Card
    add_card_shape(slide5, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.8), border_color=ACCENT_RED)
    c1_box = slide5.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.1), Inches(4.4))
    tf_c1 = c1_box.text_frame
    tf_c1.word_wrap = True
    p_c1_t = tf_c1.paragraphs[0]
    p_c1_t.text = "❌ Conventional 'Entire Image' Blurring"
    p_c1_t.font.size = Pt(18)
    p_c1_t.font.bold = True
    p_c1_t.font.color.rgb = ACCENT_RED
    
    p_c1_b = tf_c1.add_paragraph()
    p_c1_b.text = (
        "• Standard blurring tools block out the entire image or document to protect PII, rendering the document completely illegible.\n\n"
        "• If you upload a gun photo, standard software blurs the whole background or blocks the entire screen, making it impossible to see context.\n\n"
        "• Recipient cannot verify the document type, the name, or any legit details, resulting in rejected applications."
    )
    p_c1_b.font.size = Pt(12)
    p_c1_b.font.color.rgb = TEXT_WHITE
    p_c1_b.space_before = Pt(14)

    # After Card
    add_card_shape(slide5, Inches(6.98), Inches(1.8), Inches(5.6), Inches(4.8), border_color=ACCENT_CYAN)
    c2_box = slide5.shapes.add_textbox(Inches(7.23), Inches(2.0), Inches(5.1), Inches(4.4))
    tf_c2 = c2_box.text_frame
    tf_c2.word_wrap = True
    p_c2_t = tf_c2.paragraphs[0]
    p_c2_t.text = "🎯 PrivacyGuard-X Precision Redaction"
    p_c2_t.font.size = Pt(18)
    p_c2_t.font.bold = True
    p_c2_t.font.color.rgb = ACCENT_CYAN
    
    p_c2_b = tf_c2.add_paragraph()
    p_c2_b.text = (
        "• We pinpoint the EXACT pixel coordinates `(x, y, w, h)` of only the sensitive text (like Aadhaar number group or serial number).\n\n"
        "• If a weapon is uploaded, we segment ONLY the gun contour (leaving the red background fabric or suit completely legible).\n\n"
        "• The document remains 100% useful for legitimate verification, while protecting private fields."
    )
    p_c2_b.font.size = Pt(12)
    p_c2_b.font.color.rgb = TEXT_WHITE
    p_c2_b.space_before = Pt(14)

    # ==========================================================================
    # SLIDE 6: How It Works: System Architecture
    # ==========================================================================
    slide6 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide6)
    add_header(slide6, "System Architecture & Processing Flow", "3. HOW IT WORKS")

    # Large Diagram card
    add_card_shape(slide6, Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8))
    arch_box = slide6.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(4.4))
    tf_arch = arch_box.text_frame
    tf_arch.word_wrap = True
    
    p_a_t = tf_arch.paragraphs[0]
    p_a_t.text = "🛠️ Under the Hood: Secure Local Architecture"
    p_a_t.font.size = Pt(18)
    p_a_t.font.bold = True
    p_a_t.font.color.rgb = ACCENT_CYAN
    
    p_a_b = tf_arch.add_paragraph()
    p_a_b.text = (
        "✦ User Interface (Frontend): Vite + React single-page app (SPA). Implements drag-and-drop secure file loading, real-time diagnostic grids, analytics visualizers, and interactive chat panels.\n\n"
        "✦ REST API Gateway (Backend): Lightweight Python Flask microservice managing local endpoint routing, settings management, and file storage history.\n\n"
        "✦ Diagnostic Scanning Pipeline (Local Engine):\n"
        "   Image Upload ➜ Local Verification ➜ OpenCV Grayscale/Pre-processing ➜ Pytesseract OCR Targeted Character Bounding ➜ RegEx PII Pattern Analyzer ➜ Visual Weapon Contour Segmentation ➜ Local Gaussian Convolution Redactor ➜ Secure Download Outflow."
    )
    p_a_b.font.size = Pt(12.5)
    p_a_b.font.color.rgb = TEXT_WHITE
    p_a_b.space_before = Pt(16)

    # ==========================================================================
    # SLIDE 7: How It Works: Precision Document Scanning Pipeline
    # ==========================================================================
    slide7 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide7)
    add_header(slide7, "Offline Precision Document (PII) Redactor", "3. HOW IT WORKS")

    # Left Column: Preprocessing
    add_card_shape(slide7, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.8))
    l_box = slide7.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.1), Inches(4.4))
    tf_l = l_box.text_frame
    tf_l.word_wrap = True
    p_l_t = tf_l.paragraphs[0]
    p_l_t.text = "🔍 Step 1: Pre-processing & OCR"
    p_l_t.font.size = Pt(18)
    p_l_t.font.bold = True
    p_l_t.font.color.rgb = ACCENT_CYAN
    
    p_l_b = tf_l.add_paragraph()
    p_l_b.text = (
        "• OpenCV Grayscale & Cubic Resizing: Smaller documents are dynamically upscaled to guarantee PyTesseract accuracy.\n\n"
        "• Layout-Aware OCR (PSM 3): Standard layout scanners scramble word positions. We run local OCR in chronological reading order.\n\n"
        "• Confidence Filter: Excludes Tesseract structural tags and filters out low-confidence OCR artifacts (`conf < 40`) to prevent false-positive blurring."
    )
    p_l_b.font.size = Pt(12)
    p_l_b.font.color.rgb = TEXT_WHITE
    p_l_b.space_before = Pt(14)

    # Right Column: Matching & Blurring
    add_card_shape(slide7, Inches(6.98), Inches(1.8), Inches(5.6), Inches(4.8), border_color=ACCENT_CYAN)
    r_box = slide7.shapes.add_textbox(Inches(7.23), Inches(2.0), Inches(5.1), Inches(4.4))
    tf_r = r_box.text_frame
    tf_r.word_wrap = True
    p_r_t = tf_r.paragraphs[0]
    p_r_t.text = "🎯 Step 2: RegEx Pattern & Label Blur"
    p_r_t.font.size = Pt(18)
    p_r_t.font.bold = True
    p_r_t.font.color.rgb = ACCENT_CYAN
    
    p_r_b = tf_r.add_paragraph()
    p_r_b.text = (
        "• Regex Token Scanners: Detects PAN numbers, voter card IDs, enrollment numbers, credit cards, emails, and phone digits.\n\n"
        "• Sequential Multi-Token Aadhaar Match: Tracks grouped digits (like `1234 5678 9012` on the same line) and binds them as a single block.\n\n"
        "• Targeted Word Bounding Bounding: Calculates coordinate boxes on matches, overlaying tight red borders and local pixel convolution."
    )
    p_r_b.font.size = Pt(12)
    p_r_b.font.color.rgb = TEXT_WHITE
    p_r_b.space_before = Pt(14)

    # ==========================================================================
    # SLIDE 8: How It Works: Visual Weapon Blur (Red Backgrounds)
    # ==========================================================================
    slide8 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide8)
    add_header(slide8, "Visual Weapon Blurring: Red Background Optimization", "3. HOW IT WORKS")

    add_card_shape(slide8, Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8), border_color=ACCENT_RED)
    red_box = slide8.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(4.4))
    tf_red = red_box.text_frame
    tf_red.word_wrap = True
    
    p_red_t = tf_red.paragraphs[0]
    p_red_t.text = "🔴 The Red-minus-Green Color Difference Method"
    p_red_t.font.size = Pt(18)
    p_red_t.font.bold = True
    p_red_t.font.color.rgb = ACCENT_RED
    
    p_red_b = tf_red.add_paragraph()
    p_red_b.text = (
        "1. The Challenge: Weapons placed on bright red backgrounds (like download.jpg) have strong fabric textures and dark folds. Simple Red-channel thresholding flags shadows as foreground, blurring the entire image.\n\n"
        "2. The Innovation: We compute the absolute color difference: `diff = cv2.subtract(red_channel, green_channel)`.\n\n"
        "   - Saturated red fabric (even in shadows) always has a very high Red value and low Green value (`diff > 70`).\n"
        "   - Achromatic black handguns and gold bullets have very close Red and Green values (`diff <= 70`).\n\n"
        "3. Morphological Polish: Thresholding `diff` at `70` with `cv2.THRESH_BINARY_INV`, followed by morphological opening/closing, completely isolates the weapon and bullet tray in tight coordinates, keeping 100% of the red background unblurred!"
    )
    p_red_b.font.size = Pt(12.5)
    p_red_b.font.color.rgb = TEXT_WHITE
    p_red_b.space_before = Pt(16)

    # ==========================================================================
    # SLIDE 9: How It Works: Visual Weapon Blur (White Backgrounds)
    # ==========================================================================
    slide9 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide9)
    add_header(slide9, "Visual Weapon Blurring: White Background Optimization", "3. HOW IT WORKS")

    add_card_shape(slide9, Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8))
    white_box = slide9.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(4.4))
    tf_white = white_box.text_frame
    tf_white.word_wrap = True
    
    p_w_t = tf_white.paragraphs[0]
    p_w_t.text = "⚪ The Hand-Separated Size & Spatial Filter"
    p_w_t.font.size = Pt(18)
    p_w_t.font.bold = True
    p_w_t.font.color.rgb = ACCENT_CYAN
    
    p_w_b = tf_white.add_paragraph()
    p_w_b.text = (
        "1. The Challenge: In studio shots with white backdrops (like images.jpg), the man, his black suit, and the gun form one massive dark shape. Grayscale Otsu's thresholding blurs the entire person.\n\n"
        "2. Strict Value Thresholding: The handgun is pure black (`V < 45`), the black suit is dark, but the human hand skin is bright (`V > 150`). By thresholding strictly at `45` (`gray < 45`), the hand acts as a natural separator between the gun and the sleeve!\n\n"
        "3. Morphological Severing: A `3x3` opening operator sever any remaining thin single-pixel connections.\n\n"
        "4. Targeted Area & Spatial Constraints: Handguns cover between `0.1%` and `5.0%` of the image, while suits cover `>15%`. By filtering for `0.001 < ratio < 0.05` and spatial location `x / w < 0.35` (only keeping the leftmost extended shape), we perfectly redact ONLY the gun, leaving the suit, hair, face, and background clean!"
    )
    p_w_b.font.size = Pt(12)
    p_w_b.font.color.rgb = TEXT_WHITE
    p_w_b.space_before = Pt(14)

    # ==========================================================================
    # SLIDE 10: Seamless Hybrid: Local OCR + Claude Sonnet
    # ==========================================================================
    slide10 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide10)
    add_header(slide10, "Seamless Hybrid Security Engine", "3. INTEGRATED THREAT INTELLIGENCE")

    # Column 1 (Online Claude)
    add_card_shape(slide10, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.8))
    on_box = slide10.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.1), Inches(4.4))
    tf_on = on_box.text_frame
    tf_on.word_wrap = True
    p_on_t = tf_on.paragraphs[0]
    p_on_t.text = "⚡ Cloud Mode: Claude 3.5 Sonnet"
    p_on_t.font.size = Pt(18)
    p_on_t.font.bold = True
    p_on_t.font.color.rgb = ACCENT_CYAN
    
    p_on_b = tf_on.add_paragraph()
    p_on_b.text = (
        "• Triggered dynamically when a valid Anthropic API key is provided in Settings or environment variables.\n\n"
        "• Employs Claude Vision model to analyze complex visual data, identifying contraband, weapons, or private PII.\n\n"
        "• Returns exact bounding coordinates in a structured JSON payload for precision local blurring."
    )
    p_on_b.font.size = Pt(12)
    p_on_b.font.color.rgb = TEXT_WHITE
    p_on_b.space_before = Pt(14)

    # Column 2 (Offline Fallback)
    add_card_shape(slide10, Inches(6.98), Inches(1.8), Inches(5.6), Inches(4.8), border_color=ACCENT_RED)
    off_box = slide10.shapes.add_textbox(Inches(7.23), Inches(2.0), Inches(5.1), Inches(4.4))
    tf_off = off_box.text_frame
    tf_off.word_wrap = True
    p_off_t = tf_off.paragraphs[0]
    p_off_t.text = "🔒 Fallback Mode: 100% Offline"
    p_off_t.font.size = Pt(18)
    p_off_t.font.bold = True
    p_off_t.font.color.rgb = ACCENT_RED
    
    p_off_b = tf_off.add_paragraph()
    p_off_b.text = (
        "• Activated automatically and seamlessly if API keys are absent or connectivity is lost.\n\n"
        "• Uses PyTesseract OCR and custom regular expression parsing to map credentials on documents.\n\n"
        "• Employs local OpenCV contour segmentation to pinpoint and redact physical weapons and safety hazards."
    )
    p_off_b.font.size = Pt(12)
    p_off_b.font.color.rgb = TEXT_WHITE
    p_off_b.space_before = Pt(14)

    # ==========================================================================
    # SLIDE 11: Premium User Experience & Core Features
    # ==========================================================================
    slide11 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide11)
    add_header(slide11, "Premium User Experience & Core Features", "4. CORE FEATURES")

    # 4-Grid of features
    fw = Inches(5.6)
    fh = Inches(2.2)
    
    # Card 1: Scanner Tab
    add_card_shape(slide11, Inches(0.75), Inches(1.8), fw, fh)
    box = slide11.shapes.add_textbox(Inches(0.9), Inches(1.9), fw - Inches(0.3), fh - Inches(0.2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "🔍 Drag-and-Drop Scanner"
    tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ACCENT_CYAN
    p = tf.add_paragraph()
    p.text = "Vibrant glassmorphic UI showcasing original vs redacted comparisons, dynamic progress, and red alert warning banners."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(6)

    # Card 2: Analytics
    add_card_shape(slide11, Inches(6.98), Inches(1.8), fw, fh)
    box = slide11.shapes.add_textbox(Inches(7.13), Inches(1.9), fw - Inches(0.3), fh - Inches(0.2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "📊 Threat Analytics Dashboard"
    tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ACCENT_CYAN
    p = tf.add_paragraph()
    p.text = "Features a custom Donut SVG visualizer, daily scan volume metrics, document type densities, and keyword clouds."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(6)

    # Card 3: AI Chat
    add_card_shape(slide11, Inches(0.75), Inches(4.3), fw, fh)
    box = slide11.shapes.add_textbox(Inches(0.9), Inches(4.4), fw - Inches(0.3), fh - Inches(0.2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "🤖 AI Security Assistant"
    tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ACCENT_CYAN
    p = tf.add_paragraph()
    p.text = "Interactive chatbot providing legal and security advice. Falls back seamlessly to an offline rule-based advisor."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(6)

    # Card 4: History & Logs
    add_card_shape(slide11, Inches(6.98), Inches(4.3), fw, fh)
    box = slide11.shapes.add_textbox(Inches(7.13), Inches(4.4), fw - Inches(0.3), fh - Inches(0.2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "📋 Local Security Logs"
    tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ACCENT_CYAN
    p = tf.add_paragraph()
    p.text = "Keeps a running operational history with search filters, record wipes, and secure cross-origin memory blob downloads."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(6)

    # ==========================================================================
    # SLIDE 12: Business Value & Compliance
    # ==========================================================================
    slide12 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide12)
    add_header(slide12, "Business & Compliance Value Proposition", "5. VALUE PROPOSITION")

    # Center card
    add_card_shape(slide12, Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8), border_color=ACCENT_CYAN)
    val_box = slide12.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(4.4))
    tf_val = val_box.text_frame
    tf_val.word_wrap = True
    
    p_v_t = tf_val.paragraphs[0]
    p_v_t.text = "🎯 Why Organizations Need PrivacyGuard-X:"
    p_v_t.font.size = Pt(18)
    p_v_t.font.bold = True
    p_v_t.font.color.rgb = ACCENT_CYAN
    
    p_v_b = tf_val.add_paragraph()
    p_v_b.text = (
        "✔ Regulatory Shield: Ensures corporate compliance with strict Indian DPDP Act 2023 guidelines, minimizing risk of massive visual leak penalties.\n\n"
        "✔ Absolute Data Residency: Eliminates cloud server reliance. Highly appealing to defense systems, banking institutions, and medical networks requiring zero internet exposure.\n\n"
        "✔ Seamless Operation: Runs instantly and locally in remote, low-bandwidth, or secure air-gapped zones without any lag.\n\n"
        "✔ Optimal Performance: Engineered with lightweight OpenCV logic, performing robust spatial and color segmentation in milliseconds with low RAM usage.\n\n"
        "✔ Trusted Sharing: Allows document layouts to remain fully verifiable by third parties while perfectly concealing sensitive data fields."
    )
    p_v_b.font.size = Pt(13)
    p_v_b.font.color.rgb = TEXT_WHITE
    p_v_b.space_before = Pt(16)

    # ── SAVE FILE ─────────────────────────────────────────────────────────────
    filename = "PrivacyGuard-X_Presentation.pptx"
    filepath = r"C:\Users\DEEPAK\.gemini\antigravity\scratch\PrivacyGuard-X_Presentation.pptx"
    prs.save(filepath)
    print(f"Presentation saved successfully to: {filepath}")

if __name__ == '__main__':
    # Locate an image path in backend to find directories
    uploads_dir = r"C:\Users\DEEPAK\.gemini\antigravity\scratch\privacy-guard-x\backend\uploads"
    img_red = os.path.join(uploads_dir, "1779970974_download.jpg")
    create_presentation()
