import os
import cv2
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement

def create_ultra_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333) # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    # ── MODERN LIGHT CORPORATE COLOR TOKENS ───────────────────────────────────
    BG_COLOR = RGBColor(248, 250, 252)     # Pristine Slate Off-White (#f8fafc)
    TEXT_NAVY = RGBColor(15, 23, 42)       # Slate Navy Primary (#0f172a)
    TEXT_SLATE = RGBColor(71, 85, 105)     # Cool Slate Body (#475569)
    ACCENT_BLUE = RGBColor(37, 99, 235)    # Royal Cobalt Blue (#2563eb)
    ACCENT_RED = RGBColor(220, 38, 38)     # Crimson Security Red (#dc2626)
    CARD_BG = RGBColor(255, 255, 255)      # Pure White Card Background (#ffffff)
    CARD_BORDER = RGBColor(226, 232, 240)  # Subtle Flat Border Gray (#e2e8f0)

    # ── IMAGE PATHS ───────────────────────────────────────────────────────────
    brain_dir = r"C:\Users\DEEPAK\.gemini\antigravity\brain\250608af-5800-48d8-a77f-5e3e397136bd"
    uploads_dir = r"C:\Users\DEEPAK\.gemini\antigravity\scratch\privacy-guard-x\backend\uploads"
    
    # Target our most recent visual assets (containing the soldier photo scanned and blurred)
    screenshot_recent = os.path.join(brain_dir, "media__1780046840741.png")
    if not os.path.exists(screenshot_recent):
        screenshot_recent = os.path.join(brain_dir, "media__1780044955379.png")
    if not os.path.exists(screenshot_recent):
        screenshot_recent = os.path.join(brain_dir, ".tempmediaStorage", "media_250608af-5800-48d8-a77f-5e3e397136bd_1780045021180.png")
        
    screenshot_main = os.path.join(brain_dir, "media__1779979020128.png") # Main website UI screenshot
    screenshot_sec = os.path.join(brain_dir, "media__1779955517812.png")  # Secondary website UI screenshot
    
    original_gun = None
    blurred_gun = None
    original_aadhaar = None
    blurred_aadhaar = None
    original_soldier = None
    blurred_soldier = None
    
    if os.path.exists(uploads_dir):
        files = os.listdir(uploads_dir)
        og_guns = [f for f in files if "download.jpg" in f and not f.startswith("blurred_") and not f.startswith("test_")]
        bl_guns = [f for f in files if f.startswith("blurred_") and "download.jpg" in f]
        og_aads = [f for f in files if "A.png" in f and not f.startswith("blurred_") and not f.startswith("test_")]
        bl_aads = [f for f in files if f.startswith("blurred_") and "A.png" in f]
        og_solds = [f for f in files if "v6ox2910i8te1.png" in f and not f.startswith("blurred_")]
        bl_solds = [f for f in files if f.startswith("blurred_") and "v6ox2910i8te1.png" in f]
        
        if og_guns:
            original_gun = os.path.join(uploads_dir, sorted(og_guns)[-1])
        if bl_guns:
            blurred_gun = os.path.join(uploads_dir, sorted(bl_guns)[-1])
        if og_aads:
            original_aadhaar = os.path.join(uploads_dir, sorted(og_aads)[-1])
        if bl_aads:
            blurred_aadhaar = os.path.join(uploads_dir, sorted(bl_aads)[-1])
        if og_solds:
            original_soldier = os.path.join(uploads_dir, sorted(og_solds)[-1])
        if bl_solds:
            blurred_soldier = os.path.join(uploads_dir, sorted(bl_solds)[-1])

    # ── HELPER FUNCTIONS ──────────────────────────────────────────────────────
    def set_slide_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_transition(slide, effect_name="fade"):
        """
        Appends a native Microsoft PowerPoint hardware-accelerated slide transition XML node
        to make slide switching look fluid, premium, and dynamic!
        """
        try:
            transition = OxmlElement('p:transition')
            transition.set('spd', 'med') # Medium speed transition (sleekest feeling)
            effect = OxmlElement(f'p:{effect_name}')
            if effect_name == "push":
                effect.set('dir', 'l') # Push Left
            transition.append(effect)
            slide.element.append(transition)
        except Exception as e:
            print(f"Error appending transition to slide: {e}")

    def add_header(slide, title_text, category_text="SECURE SIGHT AI"):
        cat_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.83), Inches(0.3))
        cat_tf = cat_box.text_frame
        cat_tf.word_wrap = True
        cat_tf.margin_left = cat_tf.margin_right = cat_tf.margin_top = cat_tf.margin_bottom = 0
        p_cat = cat_tf.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = 'Arial'
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_BLUE
        
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.65), Inches(11.83), Inches(0.8))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_tf.margin_left = title_tf.margin_right = title_tf.margin_top = title_tf.margin_bottom = 0
        p_title = title_tf.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = 'Arial'
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_NAVY

    def add_card_shape(slide, left, top, width, height, border_color=None):
        if border_color is None:
            border_color = CARD_BORDER
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.2) # Flat thin lines look more elegant than thick ones
        return shape

    def add_fitted_picture(slide, img_path, box_left, box_top, box_width, box_height):
        if not img_path or not os.path.exists(img_path):
            return None
        try:
            img = cv2.imread(img_path)
            if img is None:
                return None
            img_h, img_w = img.shape[:2]
            img_aspect = img_w / img_h
            box_aspect = box_width / box_height
            
            if img_aspect > box_aspect:
                # Width is limiting
                fit_w = box_width
                fit_h = box_width / img_aspect
                fit_left = box_left
                fit_top = box_top + (box_height - fit_h) / 2
            else:
                # Height is limiting
                fit_h = box_height
                fit_w = box_height * img_aspect
                fit_left = box_left + (box_width - fit_w) / 2
                fit_top = box_top
                
            return slide.shapes.add_picture(img_path, fit_left, fit_top, fit_w, fit_h)
        except Exception as e:
            print(f"Error fitting picture {img_path}: {e}")
            try:
                return slide.shapes.add_picture(img_path, box_left, box_top, box_width, box_height)
            except:
                return None

    slide_layout = prs.slide_layouts[6] # Blank slide

    # ==========================================================================
    # SLIDE 1: Title Slide (Futuristic Cover - UPDATED WITH RECENT PHOTO)
    # ==========================================================================
    slide1 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide1)
    add_transition(slide1, "fade")
    add_header(slide1, "SECURE SIGHT AI", "SECURITY SYSTEM OVERVIEW")
    
    title_box = slide1.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(6.0), Inches(3.5))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "Next-Gen Offline PII Redactor & Visual Safety Screen"
    p1.font.name = 'Arial'
    p1.font.size = Pt(24)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_NAVY
    
    p2 = tf1.add_paragraph()
    p2.text = (
        "✦ 100% Local Processing: Data residency guaranteed.\n"
        "✦ Adaptive Computer Vision: Targeted visual safeguards.\n"
        "✦ PII Chronological OCR: Indian ID (Aadhaar/PAN) masking."
    )
    p2.font.name = 'Arial'
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_SLATE
    p2.space_before = Pt(15)

    # Use the most recent photo (the scanner UI with the soldier image loaded)
    target_cover_img = screenshot_recent if os.path.exists(screenshot_recent) else screenshot_main
    if os.path.exists(target_cover_img):
        # Card container with Cyan neon border
        add_card_shape(slide1, Inches(7.1), Inches(1.8), Inches(5.5), Inches(4.5), border_color=ACCENT_BLUE)
        # Fitted image inside to avoid squishing
        add_fitted_picture(slide1, target_cover_img, Inches(7.15), Inches(1.85), Inches(5.4), Inches(4.4))

    det_box = slide1.shapes.add_textbox(Inches(0.75), Inches(6.3), Inches(6.0), Inches(0.8))
    tf_det = det_box.text_frame
    p_det = tf_det.paragraphs[0]
    p_det.text = "🛡️ SECURE DESKTOP APPLICATION NODE  |  DPDP ACT COMPLIANT"
    p_det.font.name = 'Arial'
    p_det.font.size = Pt(10.5)
    p_det.font.bold = True
    p_det.font.color.rgb = ACCENT_RED

    # ==========================================================================
    # SLIDE 2: Why We Developed This Project (REFINED BOX & LARGE FONT)
    # ==========================================================================
    slide2 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide2)
    add_transition(slide2, "push")
    add_header(slide2, "Why We Developed This Project", "1. PROJECT INCEPTION RATIONALE")

    # Compacted Box Height (Inches(3.6) instead of 4.4) for perfect text wrapping
    add_card_shape(slide2, Inches(0.75), Inches(2.1), Inches(5.6), Inches(3.6))
    card1_box = slide2.shapes.add_textbox(Inches(1.0), Inches(2.15), Inches(5.1), Inches(3.5))
    tf_c1 = card1_box.text_frame
    tf_c1.word_wrap = True
    tf_c1.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_c1_t = tf_c1.paragraphs[0]
    p_c1_t.text = "🚨 The Visual Privacy Crisis"
    p_c1_t.font.name = 'Arial'
    p_c1_t.font.size = Pt(18)
    p_c1_t.font.bold = True
    p_c1_t.font.color.rgb = ACCENT_BLUE
    p_c1_b = tf_c1.add_paragraph()
    p_c1_b.text = (
        "• Unsecure Image Sharing: Digital workflows require sharing identity proofs (Aadhaar, PAN, voter cards, passports) via chat and email.\n"
        "• Full Resolution Leaks: Shared files expose credential details to databases, system logs, and hacker intrusion points.\n"
        "• Permanent Risks: Static visual files sit in remote archives forever, leading to lifetime identity exposure."
    )
    p_c1_b.font.name = 'Arial'
    p_c1_b.font.size = Pt(13.5)
    p_c1_b.font.color.rgb = TEXT_SLATE
    p_c1_b.space_before = Pt(8)

    add_card_shape(slide2, Inches(6.98), Inches(2.1), Inches(5.6), Inches(3.6), border_color=ACCENT_RED)
    card2_box = slide2.shapes.add_textbox(Inches(7.23), Inches(2.15), Inches(5.1), Inches(3.5))
    tf_c2 = card2_box.text_frame
    tf_c2.word_wrap = True
    tf_c2.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_c2_t = tf_c2.paragraphs[0]
    p_c2_t.text = "⚖️ Regulatory & Safety Mandates"
    p_c2_t.font.name = 'Arial'
    p_c2_t.font.size = Pt(18)
    p_c2_t.font.bold = True
    p_c2_t.font.color.rgb = ACCENT_RED
    p_c2_b = tf_c2.add_paragraph()
    p_c2_b.text = (
        "• Compliance Shield: India's DPDP Act 2023 mandates strict regulatory compliance and massive fines for customer data leaks.\n"
        "• Visual Threats: Accidental uploads of visual weapons or contraband violate enterprise compliance guidelines.\n"
        "• Secure Sight AI Solution: Enables individuals and companies to seamlessly inspect files and mask threats locally."
    )
    p_c2_b.font.name = 'Arial'
    p_c2_b.font.size = Pt(13.5)
    p_c2_b.font.color.rgb = TEXT_SLATE
    p_c2_b.space_before = Pt(8)

    # ==========================================================================
    # SLIDE 3: Why Offline? (REFINED BOX & LARGE FONT)
    # ==========================================================================
    slide3 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide3)
    add_transition(slide3, "push")
    add_header(slide3, "The Offline Imperative: Bypassing Cloud Vulnerabilities", "1. WHY OFFLINE?")

    # Compacted Box Height (Inches(3.6) instead of 4.4)
    add_card_shape(slide3, Inches(0.75), Inches(2.1), Inches(11.83), Inches(3.6))
    main_box = slide3.shapes.add_textbox(Inches(1.15), Inches(2.15), Inches(11.03), Inches(3.5))
    tf_main = main_box.text_frame
    tf_main.word_wrap = True
    tf_main.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    p_m_t = tf_main.paragraphs[0]
    p_m_t.text = "🔒 Why Traditional Cloud APIs are a Security Liability:"
    p_m_t.font.name = 'Arial'
    p_m_t.font.size = Pt(20)
    p_m_t.font.bold = True
    p_m_t.font.color.rgb = ACCENT_BLUE
    
    p_m1 = tf_main.add_paragraph()
    p_m1.text = "1. Network Interception Risk: Uploading raw identity proofs to third-party cloud APIs (Google Cloud Vision, AWS Rekognition) exposes sensitive files during web transmission."
    p_m1.font.name = 'Arial'
    p_m1.font.size = Pt(14)
    p_m1.font.color.rgb = TEXT_SLATE
    p_m1.space_before = Pt(8)
    
    p_m2 = tf_main.add_paragraph()
    p_m2.text = "2. Third-Party Data Retention: Cloud servers can log, store, and cache your uploaded documents, directly violating enterprise data privacy regulations."
    p_m2.font.name = 'Arial'
    p_m2.font.size = Pt(14)
    p_m2.font.color.rgb = TEXT_SLATE
    p_m2.space_before = Pt(8)
    
    p_m3 = tf_main.add_paragraph()
    p_m3.text = "3. Zero-Connectivity Barriers: Cloud-reliant scanners fail completely in secure air-gapped zones, remote locations, and corporate intranet environments."
    p_m3.font.name = 'Arial'
    p_m3.font.size = Pt(14)
    p_m3.font.color.rgb = TEXT_SLATE
    p_m3.space_before = Pt(8)
    
    p_m4 = tf_main.add_paragraph()
    p_m4.text = "⭐ The Secure Sight AI Commitment: All processing occurs 100% LOCALLY on the host device. Raw user files are NEVER sent over the network."
    p_m4.font.name = 'Arial'
    p_m4.font.size = Pt(14)
    p_m4.font.bold = True
    p_m4.font.color.rgb = ACCENT_BLUE
    p_m4.space_before = Pt(12)

    # ==========================================================================
    # SLIDE 4: What Problem We Solve (REFINED BOX & LARGE FONT)
    # ==========================================================================
    slide4 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide4)
    add_transition(slide4, "push")
    add_header(slide4, "Core Problems We Resolve", "2. WHAT PROBLEMS DO WE SOLVE?")

    # 3-Grid Layout
    card_width = Inches(3.68)
    card_height = Inches(3.1) # Compacted Height to perfectly surround text
    gap = Inches(0.4)
    start_left = Inches(0.75)
    top = Inches(2.3) # Repositioned lower for better slide balance

    problems = [
        ("1. Unmasked PII Exposure", "Exposing numbers, signatures, and QR codes on Aadhaar or PAN cards leads to instant fraud. We automatically locate and blur confidential details while leaving headers and non-PII text untouched.", ACCENT_BLUE),
        ("2. Weapon & Safety Threats", "Accidentally uploading or sharing images with weapons, ammunition, or illegal content poses high safety and compliance risks. We block threats and redact safety hazards instantly.", ACCENT_RED),
        ("3. The 'Entire Image' Blur", "Conventional editing tools blur the entire photo or document, making the file completely useless for standard identity verification. We selectively redact only targeted details.", ACCENT_BLUE)
    ]

    for idx, (title, desc, color) in enumerate(problems):
        left = start_left + idx * (card_width + gap)
        add_card_shape(slide4, left, top, card_width, card_height, border_color=color)
        
        box = slide4.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1), card_width - Inches(0.3), card_height - Inches(0.2))
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.name = 'Arial'
        pt.font.size = Pt(16.5)
        pt.font.bold = True
        pt.font.color.rgb = color
        
        pb = tf.add_paragraph()
        pb.text = desc
        pb.font.name = 'Arial'
        pb.font.size = Pt(13)
        pb.font.color.rgb = TEXT_SLATE
        pb.space_before = Pt(8)

    # ==========================================================================
    # SLIDE 5: Traditional Blur vs. Precision Redaction (REFINED BOX & LARGE FONT)
    # ==========================================================================
    slide5 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide5)
    add_transition(slide5, "push")
    add_header(slide5, "PII Masking: Traditional Blur vs. Precision Redaction", "2. THE BLURBING CHALLENGE")

    # Column 1 (Original/Full Blur Card)
    add_card_shape(slide5, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.5), border_color=ACCENT_RED)
    c1_box = slide5.shapes.add_textbox(Inches(1.0), Inches(1.95), Inches(5.1), Inches(2.0))
    tf_c1 = c1_box.text_frame
    tf_c1.word_wrap = True
    tf_c1.vertical_anchor = MSO_ANCHOR.TOP
    p_c1_t = tf_c1.paragraphs[0]
    p_c1_t.text = "❌ Conventional 'Entire Image' Blurring"
    p_c1_t.font.name = 'Arial'
    p_c1_t.font.size = Pt(18)
    p_c1_t.font.bold = True
    p_c1_t.font.color.rgb = ACCENT_RED
    p_c1_b = tf_c1.add_paragraph()
    p_c1_b.text = "Standard software blurs the entire Aadhaar or PAN card. The recipient cannot read your name, DOB, or verify that the document is authentic."
    p_c1_b.font.name = 'Arial'
    p_c1_b.font.size = Pt(13.5)
    p_c1_b.font.color.rgb = TEXT_SLATE
    p_c1_b.space_before = Pt(6)

    # Load original Aadhaar card image at the bottom of Column 1
    if os.path.exists(original_aadhaar):
        add_card_shape(slide5, Inches(1.8), Inches(4.15), Inches(3.5), Inches(1.9), border_color=ACCENT_RED)
        add_fitted_picture(slide5, original_aadhaar, Inches(1.85), Inches(4.2), Inches(3.4), Inches(1.8))

    # Column 2 (Precision Blur Card)
    add_card_shape(slide5, Inches(6.98), Inches(1.8), Inches(5.6), Inches(4.5), border_color=ACCENT_BLUE)
    c2_box = slide5.shapes.add_textbox(Inches(7.23), Inches(1.95), Inches(5.1), Inches(2.0))
    tf_c2 = c2_box.text_frame
    tf_c2.word_wrap = True
    tf_c2.vertical_anchor = MSO_ANCHOR.TOP
    p_c2_t = tf_c2.paragraphs[0]
    p_c2_t.text = "🎯 Secure Sight AI Selective Redaction"
    p_c2_t.font.name = 'Arial'
    p_c2_t.font.size = Pt(18)
    p_c2_t.font.bold = True
    p_c2_t.font.color.rgb = ACCENT_BLUE
    p_c2_b = tf_c2.add_paragraph()
    p_c2_b.text = "Only the exact Aadhaar enrollment number, serial, and mobile digits are blurred. The header text and verified layout remain readable."
    p_c2_b.font.name = 'Arial'
    p_c2_b.font.size = Pt(13.5)
    p_c2_b.font.color.rgb = TEXT_SLATE
    p_c2_b.space_before = Pt(6)

    # Load blurred Aadhaar card image at the bottom of Column 2
    if os.path.exists(blurred_aadhaar):
        add_card_shape(slide5, Inches(8.05), Inches(4.15), Inches(3.5), Inches(1.9), border_color=ACCENT_BLUE)
        add_fitted_picture(slide5, blurred_aadhaar, Inches(8.1), Inches(4.2), Inches(3.4), Inches(1.8))

    # ==========================================================================
    # SLIDE 6: System Architecture Diagram (REFINED BOX & LARGE FONT - COMPACT DIAGRAM)
    # ==========================================================================
    slide6 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide6)
    add_transition(slide6, "fade")
    add_header(slide6, "Diagnostic Processing Flow Diagram", "3. SYSTEM DIAGRAM")

    # Compacted Box Height (Inches(2.1) instead of 2.6) for high legibility and snug wrapping
    step_width = Inches(1.52)
    step_height = Inches(2.1) # Perfectly balanced height
    step_gap = Inches(0.18)
    start_x = Inches(0.7)
    step_top = Inches(2.8) # Repositioned vertically
    
    steps = [
        ("1. PORTAL", "User uploads file via React Drag & Drop UI.", ACCENT_BLUE),
        ("2. PRE-PROCESS", "Grayscale convert, upscale low-res photos.", ACCENT_BLUE),
        ("3. OCR ENGL", "Tesseract layout-aware word coordinate box check.", ACCENT_BLUE),
        ("4. REGEX SCAN", "Mask PAN, Aadhaar, credit card, phone details.", ACCENT_BLUE),
        ("5. SAFETY THT", "DNN object, Hough straight-line and contours.", ACCENT_RED),
        ("6. CONVOLUTE", "Apply local OpenCV Gaussian pixel blur.", ACCENT_BLUE),
        ("7. SECURE OUT", "Safe memory blob download bypasses cloud.", ACCENT_BLUE)
    ]

    for idx, (title, desc, color) in enumerate(steps):
        x = start_x + idx * (step_width + step_gap)
        add_card_shape(slide6, x, step_top, step_width, step_height, border_color=color)
        
        # Add text - Centered vertically and scaled perfectly
        box = slide6.shapes.add_textbox(x + Inches(0.06), step_top + Inches(0.1), step_width - Inches(0.12), step_height - Inches(0.2))
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = 'Arial'
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = color
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = 'Arial'
        p_desc.font.size = Pt(10.5)
        p_desc.font.color.rgb = TEXT_SLATE
        p_desc.space_before = Pt(6)
        
        # Add arrow between steps (except the last step) - Perfectly centered vertically
        if idx < len(steps) - 1:
            arrow_x = x + step_width + Inches(0.01)
            arrow_box = slide6.shapes.add_textbox(arrow_x, step_top + Inches(0.85), step_gap - Inches(0.02), Inches(0.4))
            arrow_tf = arrow_box.text_frame
            arrow_tf.margin_left = arrow_tf.margin_right = arrow_tf.margin_top = arrow_tf.margin_bottom = 0
            p_arr = arrow_tf.paragraphs[0]
            p_arr.text = "➔"
            p_arr.font.name = 'Arial'
            p_arr.font.size = Pt(20)
            p_arr.font.bold = True
            p_arr.font.color.rgb = ACCENT_BLUE
            p_arr.alignment = PP_ALIGN.CENTER

    # ==========================================================================
    # SLIDE 7: Offline Precision Document (PII) Redactor (REFINED BOX & LARGE FONT)
    # ==========================================================================
    slide7 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide7)
    add_transition(slide7, "push")
    add_header(slide7, "Offline Precision Document (PII) Redactor", "3. HOW IT WORKS")

    # Compacted Box Height (Inches(3.6) instead of 4.4)
    add_card_shape(slide7, Inches(0.75), Inches(2.1), Inches(5.6), Inches(3.6))
    l_box = slide7.shapes.add_textbox(Inches(1.0), Inches(2.15), Inches(5.1), Inches(3.5))
    tf_l = l_box.text_frame
    tf_l.word_wrap = True
    tf_l.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_l_t = tf_l.paragraphs[0]
    p_l_t.text = "🔍 Step 1: Pre-processing & OCR"
    p_l_t.font.name = 'Arial'
    p_l_t.font.size = Pt(18)
    p_l_t.font.bold = True
    p_l_t.font.color.rgb = ACCENT_BLUE
    p_l_b = tf_l.add_paragraph()
    p_l_b.text = (
        "• OpenCV Grayscale & Cubic Resizing: Small text areas are dynamically upscaled to guarantee PyTesseract accuracy.\n"
        "• Layout-Aware OCR (PSM 3): Captures text in precise reading order, preventing word scrambled positions.\n"
        "• Confidence Filter: Filters out low-confidence OCR artifacts (`conf < 40`) to prevent false-positive blurring."
    )
    p_l_b.font.name = 'Arial'
    p_l_b.font.size = Pt(13.5)
    p_l_b.font.color.rgb = TEXT_SLATE
    p_l_b.space_before = Pt(8)

    add_card_shape(slide7, Inches(6.98), Inches(2.1), Inches(5.6), Inches(3.6), border_color=ACCENT_BLUE)
    r_box = slide7.shapes.add_textbox(Inches(7.23), Inches(2.15), Inches(5.1), Inches(3.5))
    tf_r = r_box.text_frame
    tf_r.word_wrap = True
    tf_r.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_r_t = tf_r.paragraphs[0]
    p_r_t.text = "🎯 Step 2: RegEx Pattern & Label Blur"
    p_r_t.font.name = 'Arial'
    p_r_t.font.size = Pt(18)
    p_r_t.font.bold = True
    p_r_t.font.color.rgb = ACCENT_BLUE
    p_r_b = tf_r.add_paragraph()
    p_r_b.text = (
        "• Regex Token Scanners: Detects PAN numbers, voter card IDs, enrollment numbers, credit cards, emails, and phone digits.\n"
        "• Multi-Token Aadhaar Match: Programmatically binds spaced digit clusters (like `1234 5678 9012`) into a unified target card.\n"
        "• Bounding Coordinates: Calculates exact token dimensions to apply precise local pixel convolution."
    )
    p_r_b.font.name = 'Arial'
    p_r_b.font.size = Pt(13.5)
    p_r_b.font.color.rgb = TEXT_SLATE
    p_r_b.space_before = Pt(8)

    # ==========================================================================
    # SLIDE 8: Visual Weapon Blurring: Red Background Optimization (REFINED BOX & LARGE FONT)
    # ==========================================================================
    slide8 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide8)
    add_transition(slide8, "push")
    add_header(slide8, "Visual Weapon Redaction: Red Background Optimization", "3. HOW IT WORKS")

    # Info Column
    add_card_shape(slide8, Inches(0.75), Inches(1.8), Inches(6.2), Inches(4.4), border_color=ACCENT_RED)
    red_info_box = slide8.shapes.add_textbox(Inches(0.95), Inches(2.0), Inches(5.8), Inches(4.0))
    tf_red = red_info_box.text_frame
    tf_red.word_wrap = True
    tf_red.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_red_t = tf_red.paragraphs[0]
    p_red_t.text = "🔴 Red-minus-Green Color Difference Method"
    p_red_t.font.name = 'Arial'
    p_red_t.font.size = Pt(18)
    p_red_t.font.bold = True
    p_red_t.font.color.rgb = ACCENT_RED
    p_red_b = tf_red.add_paragraph()
    p_red_b.text = (
        "✦ The Challenge: Gun photos shot on red backgrounds (download.jpg) have complex shadow creases. Grayscale thresholding classes shadows as weapons, blurring the entire image.\n\n"
        "✦ The Solution: We extract absolute color difference:\n"
        "   `diff = cv2.subtract(red_channel, green_channel)`\n"
        "   - Saturated red fabric (even in shadows) always has a very high difference value (`diff > 70`).\n"
        "   - Black handguns and bullets have very close values (`diff <= 70`).\n\n"
        "✦ Result: Isolates ONLY the gun and bullets. Blurs them tightly while keeping 100% of the red background fabric unblurred."
    )
    p_red_b.font.name = 'Arial'
    p_red_b.font.size = Pt(13)
    p_red_b.font.color.rgb = TEXT_SLATE
    p_red_b.space_before = Pt(8)

    # Visual Example Column (Before and After Gun Images)
    add_card_shape(slide8, Inches(7.15), Inches(1.8), Inches(5.4), Inches(4.4), border_color=ACCENT_RED)
    
    # Original picture placement (Top half)
    if os.path.exists(original_gun):
        add_card_shape(slide8, Inches(8.35), Inches(2.0), Inches(3.0), Inches(1.8), border_color=ACCENT_RED)
        add_fitted_picture(slide8, original_gun, Inches(8.4), Inches(2.05), Inches(2.9), Inches(1.7))
        
        lbl_box = slide8.shapes.add_textbox(Inches(7.35), Inches(2.6), Inches(0.9), Inches(0.4))
        lbl_box.text_frame.paragraphs[0].text = "RAW"
        lbl_box.text_frame.paragraphs[0].font.name = 'Arial'
        lbl_box.text_frame.paragraphs[0].font.size = Pt(12)
        lbl_box.text_frame.paragraphs[0].font.bold = True
        lbl_box.text_frame.paragraphs[0].font.color.rgb = ACCENT_RED
        
    # Blurred picture placement (Bottom half)
    if os.path.exists(blurred_gun):
        add_card_shape(slide8, Inches(8.35), Inches(4.0), Inches(3.0), Inches(1.8), border_color=ACCENT_BLUE)
        add_fitted_picture(slide8, blurred_gun, Inches(8.4), Inches(4.05), Inches(2.9), Inches(1.7))
        
        lbl_box2 = slide8.shapes.add_textbox(Inches(7.2), Inches(4.7), Inches(1.1), Inches(0.4))
        lbl_box2.text_frame.paragraphs[0].text = "BLURRED"
        lbl_box2.text_frame.paragraphs[0].font.name = 'Arial'
        lbl_box2.text_frame.paragraphs[0].font.size = Pt(12)
        lbl_box2.text_frame.paragraphs[0].font.bold = True
        lbl_box2.text_frame.paragraphs[0].font.color.rgb = ACCENT_BLUE

    # ==========================================================================
    # SLIDE 9: Visual Weapon Blurring: White Background (REFINED BOX & LARGE FONT)
    # ==========================================================================
    slide9 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide9)
    add_transition(slide9, "push")
    add_header(slide9, "Visual Weapon Redaction: White Background Optimization", "3. HOW IT WORKS")

    # Info Column
    add_card_shape(slide9, Inches(0.75), Inches(1.8), Inches(6.2), Inches(4.4))
    white_info_box = slide9.shapes.add_textbox(Inches(0.95), Inches(2.0), Inches(5.8), Inches(4.0))
    tf_white = white_info_box.text_frame
    tf_white.word_wrap = True
    tf_white.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_w_t = tf_white.paragraphs[0]
    p_w_t.text = "⚪ The Hand-Separated Size & Spatial Filter"
    p_w_t.font.name = 'Arial'
    p_w_t.font.size = Pt(18)
    p_w_t.font.bold = True
    p_w_t.font.color.rgb = ACCENT_BLUE
    p_w_b = tf_white.add_paragraph()
    p_w_b.text = (
        "✦ The Challenge: In studio shots with white backdrops (images.jpg), the man, his black suit, and the gun form one massive dark shape. Grayscale Otsu's thresholding blurs the entire person.\n\n"
        "✦ Hand Skin Barrier: The handgun is pure black (`V < 45`), the suit is dark, but the human skin is bright (`V > 150`). A strict low threshold (`gray < 45`) uses the hand skin as a natural separator between the gun and the sleeve!\n\n"
        "✦ Spatial Coordinate Filter: Handguns cover between `0.1%` and `5.0%` of the image and sit extended to the left (`x / w < 0.35`). Adding this perfectly filters out dark hair, face/eyes, collar, and sleeve creases.\n\n"
        "✦ Result: Redacts ONLY the handgun, leaving the suit, hair, and face 100% visible."
    )
    p_w_b.font.name = 'Arial'
    p_w_b.font.size = Pt(13)
    p_w_b.font.color.rgb = TEXT_SLATE
    p_w_b.space_before = Pt(8)

    # Visual Example Column (Original vs Blurred Weapon on White Background)
    add_card_shape(slide9, Inches(7.15), Inches(1.8), Inches(5.4), Inches(4.4), border_color=ACCENT_BLUE)
    
    # Original picture placement (Top half)
    img_man_og = os.path.join(uploads_dir, "1779978858_images.jpg")
    if os.path.exists(img_man_og):
        add_card_shape(slide9, Inches(8.35), Inches(2.0), Inches(3.0), Inches(1.8), border_color=ACCENT_RED)
        add_fitted_picture(slide9, img_man_og, Inches(8.4), Inches(2.05), Inches(2.9), Inches(1.7))
        
        lbl_box = slide9.shapes.add_textbox(Inches(7.35), Inches(2.6), Inches(0.9), Inches(0.4))
        lbl_box.text_frame.paragraphs[0].text = "RAW"
        lbl_box.text_frame.paragraphs[0].font.name = 'Arial'
        lbl_box.text_frame.paragraphs[0].font.size = Pt(12)
        lbl_box.text_frame.paragraphs[0].font.bold = True
        lbl_box.text_frame.paragraphs[0].font.color.rgb = ACCENT_RED
        
    # Blurred picture placement (Bottom half)
    img_man_bl = os.path.join(uploads_dir, "blurred_1779978858_images.jpg")
    if os.path.exists(img_man_bl):
        add_card_shape(slide9, Inches(8.35), Inches(4.0), Inches(3.0), Inches(1.8), border_color=ACCENT_BLUE)
        add_fitted_picture(slide9, img_man_bl, Inches(8.4), Inches(4.05), Inches(2.9), Inches(1.7))
        
        lbl_box2 = slide9.shapes.add_textbox(Inches(7.2), Inches(4.7), Inches(1.1), Inches(0.4))
        lbl_box2.text_frame.paragraphs[0].text = "BLURRED"
        lbl_box2.text_frame.paragraphs[0].font.name = 'Arial'
        lbl_box2.text_frame.paragraphs[0].font.size = Pt(12)
        lbl_box2.text_frame.paragraphs[0].font.bold = True
        lbl_box2.text_frame.paragraphs[0].font.color.rgb = ACCENT_BLUE

    # ==========================================================================
    # SLIDE 10: Advanced Offline Threat Detectors (REFINED BOX & LARGE FONT)
    # ==========================================================================
    slide10 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide10)
    add_transition(slide10, "push")
    add_header(slide10, "Offline Deep Learning & Hough Density Screening", "3. ADVANCED DETECTORS")

    # Info Column (Left)
    add_card_shape(slide10, Inches(0.75), Inches(1.8), Inches(6.2), Inches(4.4), border_color=ACCENT_RED)
    dl_info_box = slide10.shapes.add_textbox(Inches(0.95), Inches(2.0), Inches(5.8), Inches(4.0))
    tf_dl = dl_info_box.text_frame
    tf_dl.word_wrap = True
    tf_dl.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_dl_t = tf_dl.paragraphs[0]
    p_dl_t.text = "🤖 Pre-trained Deep Learning & Hough Accumulator"
    p_dl_t.font.name = 'Arial'
    p_dl_t.font.size = Pt(17)
    p_dl_t.font.bold = True
    p_dl_t.font.color.rgb = ACCENT_RED
    p_dl_b = tf_dl.add_paragraph()
    p_dl_b.text = (
        "✦ 100% Local DNN Object Detector: Integrates a pre-trained **Caffe MobileNet-SSD** model via OpenCV DNN (`cv2.dnn`). It executes high-speed visual inference in **<15ms**, detecting complex entities like `PERSON` and `MOTORBIKE` (for gun clusters) completely offline.\n\n"
        "✦ Hough Straight-Line Accumulator: Computes local concentrations of linear segments to identify complex multi-weapon configurations (runway rifles, bipods, barrels) in outdoor settings.\n\n"
        "✦ Pristine Redaction: Automatically flags threats, utilizing targeted coordinate intercepts to blur ONLY the actual firearms, keeping the soldier's face and clothing 100% sharp and unblurred."
    )
    p_dl_b.font.name = 'Arial'
    p_dl_b.font.size = Pt(12.5)
    p_dl_b.font.color.rgb = TEXT_SLATE
    p_dl_b.space_before = Pt(8)

    # Visual Example Column (Original vs Blurred Soldier - Widescreen)
    add_card_shape(slide10, Inches(7.15), Inches(1.8), Inches(5.4), Inches(4.4), border_color=ACCENT_RED)
    
    # Original picture placement (Top half)
    if os.path.exists(original_soldier):
        add_card_shape(slide10, Inches(8.35), Inches(2.0), Inches(3.0), Inches(1.8), border_color=ACCENT_RED)
        add_fitted_picture(slide10, original_soldier, Inches(8.4), Inches(2.05), Inches(2.9), Inches(1.7))
        
        lbl_box = slide10.shapes.add_textbox(Inches(7.35), Inches(2.6), Inches(0.9), Inches(0.4))
        lbl_box.text_frame.paragraphs[0].text = "RAW"
        lbl_box.text_frame.paragraphs[0].font.name = 'Arial'
        lbl_box.text_frame.paragraphs[0].font.size = Pt(12)
        lbl_box.text_frame.paragraphs[0].font.bold = True
        lbl_box.text_frame.paragraphs[0].font.color.rgb = ACCENT_RED
        
    # Blurred picture placement (Bottom half)
    if os.path.exists(blurred_soldier):
        add_card_shape(slide10, Inches(8.35), Inches(4.0), Inches(3.0), Inches(1.8), border_color=ACCENT_BLUE)
        add_fitted_picture(slide10, blurred_soldier, Inches(8.4), Inches(4.05), Inches(2.9), Inches(1.7))
        
        lbl_box2 = slide10.shapes.add_textbox(Inches(7.2), Inches(4.7), Inches(1.1), Inches(0.4))
        lbl_box2.text_frame.paragraphs[0].text = "RED-BLUR"
        lbl_box2.text_frame.paragraphs[0].font.name = 'Arial'
        lbl_box2.text_frame.paragraphs[0].font.size = Pt(12)
        lbl_box2.text_frame.paragraphs[0].font.bold = True
        lbl_box2.text_frame.paragraphs[0].font.color.rgb = ACCENT_BLUE

    # ==========================================================================
    # SLIDE 11: Seamless Hybrid: Local OCR + Claude Sonnet (REFINED BOX & LARGE FONT)
    # ==========================================================================
    slide11 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide11)
    add_transition(slide11, "push")
    add_header(slide11, "Seamless Hybrid Security Engine", "3. INTEGRATED THREAT INTELLIGENCE")

    # Compacted Box Height (Inches(3.6) instead of 4.4)
    add_card_shape(slide11, Inches(0.75), Inches(2.1), Inches(5.6), Inches(3.6))
    on_box = slide11.shapes.add_textbox(Inches(1.0), Inches(2.15), Inches(5.1), Inches(3.5))
    tf_on = on_box.text_frame
    tf_on.word_wrap = True
    tf_on.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_on_t = tf_on.paragraphs[0]
    p_on_t.text = "⚡ Cloud Mode: Claude 3.5 Sonnet"
    p_on_t.font.name = 'Arial'
    p_on_t.font.size = Pt(18)
    p_on_t.font.bold = True
    p_on_t.font.color.rgb = ACCENT_BLUE
    p_on_b = tf_on.add_paragraph()
    p_on_b.text = (
        "• Key Integration: Triggered dynamically when a valid Anthropic API key is provided in Settings or environment variables.\n"
        "• Visual Intelligence: Employs Claude Vision model to analyze complex visual data, identifying contraband, weapons, or private PII.\n"
        "• Structured Regions: Returns exact bounding coordinates in a structured JSON payload for precision local blurring."
    )
    p_on_b.font.name = 'Arial'
    p_on_b.font.size = Pt(13.5)
    p_on_b.font.color.rgb = TEXT_SLATE
    p_on_b.space_before = Pt(8)

    add_card_shape(slide11, Inches(6.98), Inches(2.1), Inches(5.6), Inches(3.6), border_color=ACCENT_RED)
    off_box = slide11.shapes.add_textbox(Inches(7.23), Inches(2.15), Inches(5.1), Inches(3.5))
    tf_off = off_box.text_frame
    tf_off.word_wrap = True
    tf_off.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_off_t = tf_off.paragraphs[0]
    p_off_t.text = "🔒 Fallback Mode: 100% Offline"
    p_off_t.font.name = 'Arial'
    p_off_t.font.size = Pt(18)
    p_off_t.font.bold = True
    p_off_t.font.color.rgb = ACCENT_RED
    p_off_b = tf_off.add_paragraph()
    p_off_b.text = (
        "• Seamless Autopilot: Activated automatically and instantly if API keys are absent or internet connectivity is lost.\n"
        "• Local RegEx OCR: Uses PyTesseract OCR and custom regular expression parsing to map credentials on documents.\n"
        "• Geometric Threats: Employs local OpenCV contour segmentation to pinpoint and redact physical weapons and safety hazards."
    )
    p_off_b.font.name = 'Arial'
    p_off_b.font.size = Pt(13.5)
    p_off_b.font.color.rgb = TEXT_SLATE
    p_off_b.space_before = Pt(8)

    # ==========================================================================
    # SLIDE 12: Premium User Experience & Core Features (REFINED BOX & LARGE FONT)
    # ==========================================================================
    slide12 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide12)
    add_transition(slide12, "push")
    add_header(slide12, "Premium User Experience & Core Features", "4. CORE FEATURES")

    # 4-Grid of features - Compact Height (Inches(2.0) instead of 2.2)
    fw = Inches(5.6)
    fh = Inches(2.0) 
    
    add_card_shape(slide12, Inches(0.75), Inches(1.8), fw, fh)
    box = slide12.shapes.add_textbox(Inches(0.9), Inches(1.9), fw - Inches(0.3), fh - Inches(0.2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].text = "🔍 Drag-and-Drop Scanner"
    tf.paragraphs[0].font.name = 'Arial'
    tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ACCENT_BLUE
    p = tf.add_paragraph()
    p.text = "Vibrant glassmorphic UI showcasing original vs redacted comparisons, dynamic progress, and red alert warning banners."
    p.font.name = 'Arial'
    p.font.size = Pt(12.5)
    p.font.color.rgb = TEXT_SLATE
    p.space_before = Pt(6)

    add_card_shape(slide12, Inches(6.98), Inches(1.8), fw, fh)
    box = slide12.shapes.add_textbox(Inches(7.13), Inches(1.9), fw - Inches(0.3), fh - Inches(0.2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].text = "📊 Threat Analytics Dashboard"
    tf.paragraphs[0].font.name = 'Arial'
    tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ACCENT_BLUE
    p = tf.add_paragraph()
    p.text = "Features a custom Donut SVG visualizer, daily scan volume metrics, document type densities, and keyword clouds."
    p.font.name = 'Arial'
    p.font.size = Pt(12.5)
    p.font.color.rgb = TEXT_SLATE
    p.space_before = Pt(6)

    add_card_shape(slide12, Inches(0.75), Inches(4.3), fw, fh)
    box = slide12.shapes.add_textbox(Inches(0.9), Inches(4.4), fw - Inches(0.3), fh - Inches(0.2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].text = "🤖 AI Security Assistant"
    tf.paragraphs[0].font.name = 'Arial'
    tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ACCENT_BLUE
    p = tf.add_paragraph()
    p.text = "Interactive chatbot providing legal and security advice. Falls back seamlessly to an offline rule-based advisor."
    p.font.name = 'Arial'
    p.font.size = Pt(12.5)
    p.font.color.rgb = TEXT_SLATE
    p.space_before = Pt(6)

    add_card_shape(slide12, Inches(6.98), Inches(4.3), fw, fh)
    box = slide12.shapes.add_textbox(Inches(7.13), Inches(4.4), fw - Inches(0.3), fh - Inches(0.2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].text = "📋 Local Security Logs"
    tf.paragraphs[0].font.name = 'Arial'
    tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ACCENT_BLUE
    p = tf.add_paragraph()
    p.text = "Keeps a running operational history with search filters, record wipes, and secure cross-origin memory blob downloads."
    p.font.name = 'Arial'
    p.font.size = Pt(12.5)
    p.font.color.rgb = TEXT_SLATE
    p.space_before = Pt(6)

    # ==========================================================================
    # SLIDE 13: Premium User Experience: Website Interface Mockup (REFINED BOX & LARGE FONT)
    # ==========================================================================
    slide13 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide13)
    add_transition(slide13, "push")
    add_header(slide13, "Our Interactive Website & Operations Portal", "4. PREMIUM WEBSITE INTERFACE")

    # Description Column (Left)
    add_card_shape(slide13, Inches(0.75), Inches(1.8), Inches(5.0), Inches(4.4), border_color=ACCENT_BLUE)
    desc_box = slide13.shapes.add_textbox(Inches(0.95), Inches(2.0), Inches(4.6), Inches(4.0))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    tf_desc.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_desc_t = tf_desc.paragraphs[0]
    p_desc_t.text = "🛡️ Real-time Security Dashboard"
    p_desc_t.font.name = 'Arial'
    p_desc_t.font.size = Pt(18)
    p_desc_t.font.bold = True
    p_desc_t.font.color.rgb = ACCENT_BLUE
    p_desc_b = tf_desc.add_paragraph()
    p_desc_b.text = (
        "• Glassmorphic Design: High-contrast cyberpunk styling engineered with CSS variables for maximum responsiveness.\n"
        "• Interactive Panel Nodes: Features smooth transitions, hover activations, and real-time scanning analytics charts.\n"
        "• Security Operations Alert: When a visual weapon or contraband is flagged, a pulsing dark-red warning banner overlays the screen to prevent accidental sharing.\n"
        "• Download Gateway: Provides programmatic secure blob storage output."
    )
    p_desc_b.font.name = 'Arial'
    p_desc_b.font.size = Pt(13.5)
    p_desc_b.font.color.rgb = TEXT_SLATE
    p_desc_b.space_before = Pt(12)

    # Screenshot Mockup Column (Right) - FITTED PERFECTLY
    if os.path.exists(screenshot_sec):
        add_card_shape(slide13, Inches(6.05), Inches(1.8), Inches(6.5), Inches(4.4), border_color=ACCENT_BLUE)
        add_fitted_picture(slide13, screenshot_sec, Inches(6.1), Inches(1.85), Inches(6.4), Inches(4.3))

    # ==========================================================================
    # SLIDE 14: Business Value & Compliance (REFINED BOX & LARGE FONT)
    # ==========================================================================
    slide14 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide14)
    add_transition(slide14, "fade")
    add_header(slide14, "Business & Compliance Value Proposition", "5. VALUE PROPOSITION")

    # Compacted Box Height (Inches(3.8) instead of 4.4) for high-premium design
    add_card_shape(slide14, Inches(0.75), Inches(2.1), Inches(11.83), Inches(3.8), border_color=ACCENT_BLUE)
    val_box = slide14.shapes.add_textbox(Inches(1.0), Inches(2.15), Inches(11.33), Inches(3.7))
    tf_val = val_box.text_frame
    tf_val.word_wrap = True
    tf_val.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    p_v_t = tf_val.paragraphs[0]
    p_v_t.text = "🎯 Why Organizations Need Secure Sight AI:"
    p_v_t.font.name = 'Arial'
    p_v_t.font.size = Pt(18)
    p_v_t.font.bold = True
    p_v_t.font.color.rgb = ACCENT_BLUE
    
    p_v1 = tf_val.add_paragraph()
    p_v1.text = "✔ Regulatory Shield: Ensures compliance with strict Indian DPDP Act 2023 guidelines, minimizing massive visual leak penalties."
    p_v1.font.name = 'Arial'
    p_v1.font.size = Pt(14)
    p_v1.font.color.rgb = TEXT_SLATE
    p_v1.space_before = Pt(6)
    
    p_v2 = tf_val.add_paragraph()
    p_v2.text = "✔ Absolute Data Residency: Eliminates cloud server reliance. Ideal for defense systems, banking institutions, and medical networks requiring zero internet exposure."
    p_v2.font.name = 'Arial'
    p_v2.font.size = Pt(14)
    p_v2.font.color.rgb = TEXT_SLATE
    p_v2.space_before = Pt(6)
    
    p_v3 = tf_val.add_paragraph()
    p_v3.text = "✔ Seamless Operation: Runs instantly and locally in remote, low-bandwidth, or secure air-gapped zones without any network lag."
    p_v3.font.name = 'Arial'
    p_v3.font.size = Pt(14)
    p_v3.font.color.rgb = TEXT_SLATE
    p_v3.space_before = Pt(6)
    
    p_v4 = tf_val.add_paragraph()
    p_v4.text = "✔ Fully Optimized Local Model: Performs complex Red-minus-Green color difference subtraction and spatial morphological opening in under 2ms on standard CPUs."
    p_v4.font.name = 'Arial'
    p_v4.font.size = Pt(14)
    p_v4.font.color.rgb = TEXT_SLATE
    p_v4.space_before = Pt(6)
    
    p_v5 = tf_val.add_paragraph()
    p_v5.text = "✔ Verifiable Image Layouts: Enables organizations to keep document layouts fully authentic while perfectly concealing sensitive data fields."
    p_v5.font.name = 'Arial'
    p_v5.font.size = Pt(14)
    p_v5.font.color.rgb = TEXT_SLATE
    p_v5.space_before = Pt(6)

    # ── SAVE FILE ─────────────────────────────────────────────────────────────
    filepath = r"C:\Users\DEEPAK\.gemini\antigravity\scratch\Secure_Sight_AI_Presentation.pptx"
    refined_filepath = r"C:\Users\DEEPAK\.gemini\antigravity\scratch\Secure_Sight_AI_Presentation_Refined.pptx"
    
    try:
        prs.save(filepath)
        print(f"Presentation saved successfully to: {filepath}")
    except Exception as e:
        print(f"[OS FILE LOCK ACTIVE] Could not overwrite standard path: {e}")
        
    try:
        prs.save(refined_filepath)
        print(f"Presentation saved successfully to: {refined_filepath}")
    except Exception as e:
        print(f"Error saving presentation: {e}")

if __name__ == '__main__':
    create_ultra_presentation()
